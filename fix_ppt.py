from pptx import Presentation


PPT_PATH = r"F:\РАБОТА\Кутузовский 32 (Г)\K32G_updated_v2.pptx"


def set_shape_text(slide, idx, text):
    slide.shapes[idx - 1].text = text


def set_table_rows(shape, rows):
    table = shape.table
    for r, (label, value) in enumerate(rows, start=1):
        table.cell(r - 1, 0).text = label
        table.cell(r - 1, 1).text = value


prs = Presentation(PPT_PATH)

budget_dual = "Подрядный: 94,0 млн ₽\nСобственный: 127,7 млн ₽"
works_dual = "Подрядный: 28,2 млн ₽\nСобственный: 49,5 млн ₽"

# Slide 1
slide = prs.slides[0]
set_shape_text(slide, 15, "Бюджет проекта")
set_shape_text(slide, 16, budget_dual)
set_shape_text(slide, 19, "127,7 млн ₽")
set_shape_text(slide, 22, "94,0 млн ₽")

# Slide 2
slide = prs.slides[1]
set_shape_text(slide, 8, "Итог модели (бюджет проекта)")
set_shape_text(slide, 9, budget_dual)
set_shape_text(slide, 17, "Работы / ПНР")
set_shape_text(slide, 18, works_dual)
set_shape_text(slide, 74, "Распределение бюджета")
set_shape_text(slide, 76, "Подрядный: 94,0 млн ₽ / 1 554 ₽/м²\nСобственный: 127,7 млн ₽ / 2 110 ₽/м²")
set_shape_text(slide, 81, "Удельный бюджет")
set_shape_text(slide, 82, "Подрядный: 1 554 ₽ / м²\nСобственный: 2 110 ₽ / м²")

# Slide 4
slide = prs.slides[3]
set_shape_text(slide, 14, "Итого по системе (бюджет проекта)")
set_shape_text(slide, 15, budget_dual)
set_shape_text(slide, 24, works_dual)
set_shape_text(
    slide,
    29,
    "Подрядный: Обор. 46,0% • Матер. 14,9% • Работы 36,2% • Проект. 2,8%\n"
    "Собственный: Обор. 36,1% • Матер. 11,7% • Работы 49,9% • Проект. 2,2%",
)
set_shape_text(slide, 51, "486 / 853 ₽ за 1 м кабеля")
set_shape_text(slide, 52, "подрядный / собственный ресурс")
set_shape_text(slide, 54, "12 423 / 21 797 ₽ за 1 устройство")
set_shape_text(slide, 55, "монтаж / адресация")
set_shape_text(slide, 57, "")
set_shape_text(slide, 58, "")

# Slide 6
slide = prs.slides[5]
set_shape_text(slide, 12, works_dual)
set_shape_text(slide, 17, "Бюджет проекта")
set_shape_text(slide, 18, budget_dual)
set_shape_text(slide, 21, "127,7 млн ₽")
set_shape_text(slide, 24, "94,0 млн ₽")
set_shape_text(slide, 27, "1 554 / 2 110 ₽")
set_shape_text(slide, 32, "Удельный бюджет: 2 110 ₽ / м²")
set_shape_text(slide, 33, "Работы / ПНР: подрядный 28,2 млн ₽ • собственный 49,5 млн ₽")
set_shape_text(slide, 39, "Удельный бюджет: 1 554 ₽ / м²")

set_table_rows(
    slide.shapes[29],
    [
        ("База расчета по работам", "28 200 000 ₽"),
        ("ФОТ", "48 341 849 ₽"),
        ("ФОТ ОПР", "1 137 709 ₽"),
        ("Работы / ПНР собственным ресурсом", "49 479 558 ₽"),
        ("Инструмент / расходники / СИЗ", "4 834 185 ₽"),
        ("АХР и ОПР", "23 777 134 ₽"),
        ("Оборудование (без наценки 20%)", "35 800 000 ₽"),
        ("Материалы (без наценки 15%)", "11 600 000 ₽"),
        ("Проектирование", "2 209 864 ₽"),
        ("Итог бюджета проекта", "127 700 740 ₽"),
    ],
)

set_table_rows(
    slide.shapes[36],
    [
        ("Базовая стоимость работ", "28 200 000 ₽"),
        ("Корректировка работ на коэффициенты", "13 553 625 ₽"),
        ("ОПР + АХР и ОПР", "2 677 592 ₽"),
        ("Оборудование (без наценки 20%)", "35 800 000 ₽"),
        ("Материалы (без наценки 15%)", "11 600 000 ₽"),
        ("Проектирование", "2 209 864 ₽"),
        ("Итог бюджета проекта", "94 041 081 ₽"),
    ],
)

prs.save(PPT_PATH)
print(PPT_PATH)
