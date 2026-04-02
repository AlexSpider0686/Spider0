from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path.cwd()
ASSETS = ROOT / "presentation_assets"
SCREENSHOTS = ASSETS / "screenshots"
RENDERED = ASSETS / "rendered"
OUTPUT = ROOT / "ProjectCore_Product_Presentation_Reference_Style_2026-03-30.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

BG = RGBColor(244, 247, 251)
SURFACE = RGBColor(255, 255, 255)
INK = RGBColor(15, 23, 42)
MUTED = RGBColor(88, 102, 126)
LINE = RGBColor(214, 223, 234)
BLUE = RGBColor(46, 113, 255)
CYAN = RGBColor(71, 201, 235)
ORANGE = RGBColor(255, 127, 63)
GREEN = RGBColor(45, 192, 139)
NAVY = RGBColor(10, 20, 35)


def style(run, size, bold=False, color=INK, name="Aptos"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_full_bg(slide, image_path=None, color=BG):
    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), 0, 0, width=Inches(SLIDE_W), height=Inches(SLIDE_H))
    else:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(SLIDE_H))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()


def add_overlay(slide, x, y, w, h, color, transparency=0.0, rounded=True, line_color=None, line_transparency=0.0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.transparency = line_transparency
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    style(run, size, bold, color)
    return box


def add_bullets(slide, items, x, y, w, h, size=16, color=MUTED):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = f"• {item}"
        style(run, size, False, color)
    return box


def add_chip(slide, text, x, y, accent=BLUE, width=None):
    width = width or max(1.2, min(3.0, 0.115 * len(text) + 0.72))
    shape = add_overlay(slide, x, y, width, 0.36, SURFACE, transparency=0.02, rounded=True, line_color=accent, line_transparency=0.15)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    style(run, 11, True, accent)


def add_metric(slide, x, y, w, value, label, accent=BLUE):
    add_overlay(slide, x, y, w, 1.1, SURFACE, transparency=0.0, rounded=True, line_color=LINE, line_transparency=0.0)
    add_overlay(slide, x, y, 0.08, 1.1, accent, transparency=0.0, rounded=False)
    add_text(slide, value, x + 0.22, y + 0.14, w - 0.3, 0.34, 22, True, INK)
    add_text(slide, label, x + 0.22, y + 0.63, w - 0.3, 0.22, 10.5, False, MUTED)


def add_device_frame(slide, image_path, x, y, w, h, dark=False):
    shadow = add_overlay(slide, x + 0.08, y + 0.12, w, h, NAVY if dark else RGBColor(186, 197, 212), transparency=0.65 if dark else 0.78)
    shadow.line.fill.background()
    frame = add_overlay(slide, x, y, w, h, NAVY if dark else SURFACE, transparency=0.0, rounded=True, line_color=RGBColor(42, 55, 78) if dark else LINE)
    screen = add_overlay(slide, x + 0.13, y + 0.18, w - 0.26, h - 0.31, RGBColor(0, 0, 0) if dark else RGBColor(248, 250, 252), transparency=0.0, rounded=True)
    screen.line.fill.background()
    slide.shapes.add_picture(str(image_path), Inches(x + 0.16), Inches(y + 0.21), width=Inches(w - 0.32), height=Inches(h - 0.37))
    return frame


def add_footer(slide, n):
    add_text(slide, "Project.Core™ / reference-inspired deck", 0.55, 7.03, 3.3, 0.18, 9.5, False, MUTED)
    add_text(slide, str(n), 12.25, 7.03, 0.4, 0.18, 9.5, True, MUTED, align=PP_ALIGN.RIGHT)


def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    bg_cover = RENDERED / "bg-cover.jpg"
    bg_site = RENDERED / "bg-site.jpg"
    bg_platform = RENDERED / "bg-platform.jpg"
    bg_split = RENDERED / "bg-split.jpg"
    bg_tariffs = RENDERED / "bg-tariffs.jpg"
    bg_roadmap = RENDERED / "bg-roadmap.jpg"

    # 1 cover
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, bg_cover)
    add_overlay(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY, transparency=0.38, rounded=False)
    add_overlay(slide, 0.65, 0.58, 5.65, 6.25, NAVY, transparency=0.18, rounded=True, line_color=RGBColor(255, 255, 255), line_transparency=0.75)
    add_chip(slide, "Inspired by modern SaaS decks", 0.9, 0.78, accent=CYAN, width=2.45)
    add_text(slide, "Project.Core™", 0.9, 1.32, 4.4, 0.55, 30, True, SURFACE)
    add_text(slide, "Сайт и платформа\nпредварительной бюджетной оценки\nсистем безопасности", 0.9, 1.86, 5.05, 1.55, 25, True, SURFACE)
    add_text(slide, "Новая версия презентации в более современной SaaS-стилистике: крупные product-shot’ы, чистая сетка, акцентные метрики и ясное деление на продуктовые контуры.", 0.9, 3.7, 4.75, 0.9, 13, False, RGBColor(226, 232, 240))
    add_chip(slide, "6 систем", 0.92, 5.02, accent=ORANGE)
    add_chip(slide, "5–10 минут", 2.18, 5.02, accent=CYAN)
    add_chip(slide, "Сайт + платформа", 3.82, 5.02, accent=GREEN, width=1.9)
    add_device_frame(slide, SCREENSHOTS / "site-hero.png", 6.9, 0.72, 5.75, 6.02, dark=True)

    # 2 product overview
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, None, BG)
    add_text(slide, "Продукт в одном экране", 0.68, 0.62, 4.2, 0.4, 28, True, INK)
    add_text(slide, "Project.Core™ состоит из двух визуально разных, но логически связанных частей: сайта для презентации ценности и платформы для расчёта бюджета.", 0.68, 1.06, 7.4, 0.46, 12.5, False, MUTED)
    add_metric(slide, 0.7, 1.82, 2.45, "6", "подсистем в одном расчёте", BLUE)
    add_metric(slide, 3.32, 1.82, 2.45, "85+", "субъектов РФ в модели", CYAN)
    add_metric(slide, 5.94, 1.82, 2.45, "AI", "аудит цен и рисков", ORANGE)
    add_metric(slide, 8.56, 1.82, 2.45, "B2B", "корпоративный и рыночный контур", GREEN)
    add_overlay(slide, 0.7, 3.15, 5.85, 3.55, SURFACE, line_color=LINE)
    add_overlay(slide, 6.78, 3.15, 5.85, 3.55, SURFACE, line_color=LINE)
    add_text(slide, "Сайт", 0.95, 3.42, 1.0, 0.2, 20, True, INK)
    add_text(slide, "Показывает продукт, преимущества, AI-логику и legal-контур.", 0.95, 3.72, 3.9, 0.35, 12, False, MUTED)
    add_text(slide, "Платформа", 7.02, 3.42, 1.8, 0.2, 20, True, INK)
    add_text(slide, "Считает объект, системы, проектирование, бюджет и риск-контур.", 7.02, 3.72, 4.1, 0.35, 12, False, MUTED)
    add_device_frame(slide, SCREENSHOTS / "site-hero.png", 0.96, 4.15, 5.0, 2.15, dark=False)
    add_device_frame(slide, SCREENSHOTS / "platform-object-view.png", 7.04, 4.15, 4.92, 2.15, dark=False)
    add_footer(slide, 2)

    # 3 website
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, bg_site)
    add_overlay(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY, transparency=0.46, rounded=False)
    add_text(slide, "Сайт продукта", 0.72, 0.58, 3.0, 0.34, 28, True, SURFACE)
    add_text(slide, "Современная логика подачи для B2B software: сильный hero-блок, доказательные метрики, разбор преимуществ и понятный переход к демо-сценарию.", 0.72, 1.02, 5.45, 0.58, 12.5, False, RGBColor(228, 234, 242))
    add_bullets(slide, [
        "Hero-экран с коротким УТП и CTA.",
        "Блок сравнения с альтернативами рынка.",
        "Выделенный AI-контур и объяснение методологии.",
        "Подробная страница «О системе» и legal-раздел."
    ], 0.82, 1.92, 4.6, 2.0, size=15, color=RGBColor(233, 238, 246))
    add_device_frame(slide, SCREENSHOTS / "site-hero.png", 6.15, 0.78, 6.28, 2.7, dark=True)
    add_device_frame(slide, SCREENSHOTS / "site-comparison.png", 6.15, 3.78, 3.02, 2.84, dark=True)
    add_device_frame(slide, SCREENSHOTS / "site-ai-engine.png", 9.42, 3.78, 3.02, 2.84, dark=True)
    add_footer(slide, 3)

    # 4 platform
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, bg_platform)
    add_overlay(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY, transparency=0.48, rounded=False)
    add_text(slide, "Платформа", 0.72, 0.58, 2.4, 0.34, 28, True, SURFACE)
    add_text(slide, "Визуально платформа ближе к хорошему SaaS-dashboard: шаги процесса сверху, крупные рабочие карточки, видимые экспортные сценарии и объяснимость результата.", 0.72, 1.0, 6.0, 0.56, 12.5, False, RGBColor(228, 234, 242))
    add_device_frame(slide, SCREENSHOTS / "platform-object-view.png", 0.72, 1.84, 6.15, 4.95, dark=True)
    add_device_frame(slide, SCREENSHOTS / "platform-systems-view.png", 7.02, 1.84, 5.58, 3.06, dark=True)
    add_device_frame(slide, SCREENSHOTS / "platform-budget-view.png", 7.02, 5.1, 2.67, 1.56, dark=True)
    add_device_frame(slide, SCREENSHOTS / "platform-risks-view.png", 9.93, 5.1, 2.67, 1.56, dark=True)
    add_footer(slide, 4)

    # 5 functions
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, None, BG)
    add_text(slide, "Ключевой функционал платформы", 0.68, 0.62, 5.2, 0.4, 28, True, INK)
    add_text(slide, "Слайд построен по логике product overview из современных deck-шаблонов: меньше описательности, больше функциональных окон и явных value points.", 0.68, 1.05, 7.0, 0.4, 12.5, False, MUTED)
    cards = [
        (0.7, 1.82, 3.95, 2.08, "Объект и зонирование", "Тип объекта, площади, этажность, регион, шаблоны зон и AI-обследование."),
        (4.7, 1.82, 3.95, 2.08, "Системы и спецификация", "Состав систем, вендоры, PDF APS, ключевое оборудование и рассчитанная спецификация."),
        (8.7, 1.82, 3.95, 2.08, "Бюджет и коэффициенты", "Условия расчёта, защитные механизмы, логика работ и агрегированная стоимость."),
        (0.7, 4.14, 3.95, 2.08, "AI-логика", "AI-аудит цен, AI-обследование, AI-техническое решение и AI-риски проекта."),
        (4.7, 4.14, 3.95, 2.08, "Экспорт", "ТКП, план проекта, Excel-выгрузки и презентационные материалы."),
        (8.7, 4.14, 3.95, 2.08, "Explainability", "Пользователь видит не только итог, но и происхождение суммы."),
    ]
    accents = [BLUE, CYAN, ORANGE, GREEN, BLUE, CYAN]
    for idx, card in enumerate(cards):
        x, y, w, h, title, desc = card
        add_overlay(slide, x, y, w, h, SURFACE, line_color=LINE)
        add_overlay(slide, x, y, 0.08, h, accents[idx], rounded=False)
        add_text(slide, title, x + 0.18, y + 0.22, w - 0.35, 0.26, 17, True, INK)
        add_text(slide, desc, x + 0.18, y + 0.62, w - 0.35, 1.0, 12, False, MUTED)
    add_footer(slide, 5)

    # 6 roadmap
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, bg_roadmap)
    add_overlay(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY, transparency=0.48, rounded=False)
    add_text(slide, "Планы по доработкам", 0.72, 0.6, 4.8, 0.34, 28, True, SURFACE)
    add_text(slide, "После завершения написания основного кода развитие продукта идёт по трём направлениям: технологическому, нормативному и корпоративному.", 0.72, 1.03, 7.4, 0.42, 12.5, False, RGBColor(228, 234, 242))
    steps = [
        ("01", "Завершение основной разработки", "Финишная стабилизация сайта и платформы как единого продукта."),
        ("02", "Переход на отечественное ПО", "Запланирован после завершения основного контура кода."),
        ("03", "Правовая регистрация", "Товарный знак и регистрация продукта в соответствии с законодательством РФ."),
        ("04", "Корпоративная secure-ветка", "Отдельная версия для ООО «СТК» (ПАО Сбер) с соблюдением стандартов кибербезопасности РФ и ПАО Сбер."),
    ]
    y = 1.82
    for idx, (n, title, desc) in enumerate(steps):
        add_overlay(slide, 0.84, y, 11.7, 1.12, SURFACE, transparency=0.06, line_color=RGBColor(255, 255, 255), line_transparency=0.62)
        add_text(slide, n, 1.08, y + 0.23, 0.5, 0.3, 18, True, CYAN)
        add_text(slide, title, 1.72, y + 0.18, 3.9, 0.28, 17, True, SURFACE)
        add_text(slide, desc, 5.05, y + 0.18, 6.9, 0.5, 12, False, RGBColor(226, 232, 240))
        y += 1.28
    add_footer(slide, 6)

    # 7 split products
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, bg_split)
    add_overlay(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY, transparency=0.5, rounded=False)
    add_text(slide, "Два продукта", 0.72, 0.58, 3.2, 0.34, 28, True, SURFACE)
    add_text(slide, "Аналогично сильным B2B deck’ам, продуктовая стратегия упакована в два отдельных коммерческих контура.", 0.72, 1.02, 6.5, 0.4, 12.5, False, RGBColor(228, 234, 242))
    add_overlay(slide, 0.72, 1.76, 5.92, 4.95, NAVY, transparency=0.16, line_color=ORANGE, line_transparency=0.15)
    add_overlay(slide, 6.7, 1.76, 5.92, 4.95, NAVY, transparency=0.16, line_color=CYAN, line_transparency=0.15)
    add_chip(slide, "Контур 1", 0.98, 2.02, accent=ORANGE, width=1.22)
    add_text(slide, "Версия для ООО «СТК»\n(под Генеральное соглашение с ПАО Сбер)", 0.98, 2.45, 4.6, 0.9, 20, True, SURFACE)
    add_bullets(slide, [
        "Приведение решения в соответствие требованиям стандартов по кибербезопасности РФ и ПАО Сбер.",
        "Отдельный корпоративный контур поставки и сопровождения.",
        "Фокус на compliance, безопасную эксплуатацию и внутренние требования заказчика."
    ], 0.98, 3.62, 4.82, 2.2, size=14.5, color=RGBColor(231, 236, 244))
    add_chip(slide, "Контур 2", 6.96, 2.02, accent=CYAN, width=1.22)
    add_text(slide, "Внешний рынок:\nпродажа доступа к платформе", 6.96, 2.45, 4.4, 0.9, 20, True, SURFACE)
    add_bullets(slide, [
        "SaaS-модель с тарифными планами и разным доступом к функционалу.",
        "Браузерный доступ к платформе как коммерческому продукту.",
        "Абонентская плата и масштабирование на несколько сегментов заказчиков."
    ], 6.96, 3.62, 4.82, 2.2, size=14.5, color=RGBColor(231, 236, 244))
    add_footer(slide, 7)

    # 8 tariffs
    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, bg_tariffs)
    add_overlay(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY, transparency=0.52, rounded=False)
    add_text(slide, "Коммерческая модель внешнего рынка", 0.72, 0.58, 6.4, 0.34, 28, True, SURFACE)
    add_text(slide, "По аналогии с лучшими software-презентациями, тарифы показаны как продуктовые уровни, а не как перегруженная таблица прайс-листа.", 0.72, 1.02, 7.8, 0.42, 12.5, False, RGBColor(228, 234, 242))
    plans = [
        ("Start", "от 49 000 ₽/мес", ["Базовый пресейл", "Ограниченный доступ", "Стандартные выгрузки"], BLUE),
        ("Pro", "от 119 000 ₽/мес", ["Полный расчёт", "AI-модули", "Расширенные выходные материалы"], ORANGE),
        ("Enterprise", "индивидуально", ["Корпоративная конфигурация", "Интеграции и SLA", "Отдельный контур сопровождения"], GREEN),
    ]
    xs = [0.86, 4.46, 8.06]
    for idx, (name, price, feats, accent) in enumerate(plans):
        add_overlay(slide, xs[idx], 1.86, 3.02, 4.95, SURFACE, transparency=0.02, line_color=accent, line_transparency=0.0)
        add_text(slide, name, xs[idx] + 0.2, 2.18, 2.2, 0.25, 20, True, INK)
        add_text(slide, price, xs[idx] + 0.2, 2.62, 2.4, 0.28, 16, True, accent)
        add_bullets(slide, feats, xs[idx] + 0.2, 3.25, 2.55, 1.8, size=13.4, color=MUTED)
        add_chip(slide, "Абонентская плата", xs[idx] + 0.2, 5.98, accent=accent, width=1.7)
    add_text(slide, "Для версии ООО «СТК» (ПАО Сбер) предполагается отдельная договорная модель вне стандартной SaaS-тарификации.", 0.86, 6.98, 11.6, 0.2, 10.5, False, RGBColor(214, 221, 230))
    add_footer(slide, 8)

    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    output = build()
    print(output)
