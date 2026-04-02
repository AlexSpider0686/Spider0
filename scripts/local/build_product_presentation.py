from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path.cwd()
ASSETS = ROOT / "presentation_assets"
SCREENSHOTS = ASSETS / "screenshots"
RENDERED = ASSETS / "rendered"
PUBLIC = ROOT / "public" / "assets"

OUTPUT = ROOT / "ProjectCore_Product_Presentation_2026-03-30.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
PX_W = 1600
PX_H = 900

COLORS = {
    "navy": RGBColor(11, 23, 41),
    "navy_soft": RGBColor(21, 37, 62),
    "cyan": RGBColor(99, 212, 247),
    "orange": RGBColor(255, 138, 61),
    "green": RGBColor(138, 209, 133),
    "white": RGBColor(255, 255, 255),
    "mist": RGBColor(227, 236, 245),
    "slate": RGBColor(154, 173, 194),
    "ink": RGBColor(28, 44, 69),
}


def ensure_dirs():
    RENDERED.mkdir(parents=True, exist_ok=True)


def font(size, bold=False, color=COLORS["white"], name="Aptos"):
    return {"size": Pt(size), "bold": bold, "color": color, "name": name}


def cover_resize(src: Path, dst: Path, size=(PX_W, PX_H), darken=0):
    image = Image.open(src).convert("RGB")
    src_ratio = image.width / image.height
    dst_ratio = size[0] / size[1]
    if src_ratio > dst_ratio:
        new_height = size[1]
        new_width = int(new_height * src_ratio)
    else:
        new_width = size[0]
        new_height = int(new_width / src_ratio)
    image = image.resize((new_width, new_height), Image.LANCZOS)
    left = (new_width - size[0]) // 2
    top = (new_height - size[1]) // 2
    image = image.crop((left, top, left + size[0], top + size[1]))
    if darken > 0:
        overlay = Image.new("RGBA", size, (5, 12, 24, darken))
        image = Image.alpha_composite(image.convert("RGBA"), overlay).convert("RGB")
    image.save(dst, quality=94)
    return dst


def top_crop(src: Path, dst: Path, size=(960, 540), top_share=0.0):
    image = Image.open(src).convert("RGB")
    src_ratio = image.width / image.height
    dst_ratio = size[0] / size[1]
    if src_ratio > dst_ratio:
        new_height = size[1]
        new_width = int(new_height * src_ratio)
    else:
        new_width = size[0]
        new_height = int(new_width / src_ratio)
    image = image.resize((new_width, new_height), Image.LANCZOS)
    max_top = max(new_height - size[1], 0)
    top = int(max_top * top_share)
    left = max((new_width - size[0]) // 2, 0)
    image = image.crop((left, top, left + size[0], top + size[1]))
    image.save(dst, quality=94)
    return dst


def build_assets():
    ensure_dirs()
    backgrounds = {
        "cover": cover_resize(PUBLIC / "background" / "development-lab.jpg", RENDERED / "bg-cover.jpg", darken=20),
        "overview": cover_resize(PUBLIC / "metrics" / "systems-overview.jpg", RENDERED / "bg-overview.jpg", darken=35),
        "site": cover_resize(PUBLIC / "object-types" / "public.jpg", RENDERED / "bg-site.jpg", darken=55),
        "site_blocks": cover_resize(PUBLIC / "metrics" / "multi-device.jpg", RENDERED / "bg-site-blocks.jpg", darken=60),
        "platform": cover_resize(PUBLIC / "object-types" / "production.jpg", RENDERED / "bg-platform.jpg", darken=60),
        "platform2": cover_resize(PUBLIC / "metrics" / "risk-guard.jpg", RENDERED / "bg-platform2.jpg", darken=65),
        "roadmap": cover_resize(PUBLIC / "object-types" / "transport.jpg", RENDERED / "bg-roadmap.jpg", darken=60),
        "split": cover_resize(PUBLIC / "object-types" / "warehouse.jpg", RENDERED / "bg-split.jpg", darken=65),
        "tariffs": cover_resize(PUBLIC / "object-types" / "energy.jpg", RENDERED / "bg-tariffs.jpg", darken=60),
        "about_preview": top_crop(SCREENSHOTS / "site-about-system.png", RENDERED / "site-about-preview.jpg", top_share=0.0),
    }
    return backgrounds


def set_run_style(run, style):
    run.font.size = style["size"]
    run.font.bold = style["bold"]
    run.font.color.rgb = style["color"]
    run.font.name = style["name"]


def add_bg(slide, image_path: Path):
    slide.shapes.add_picture(str(image_path), 0, 0, width=Inches(SLIDE_W), height=Inches(SLIDE_H))


def add_overlay(slide, color, transparency=0.2, x=0, y=0, w=SLIDE_W, h=SLIDE_H):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    fill.transparency = transparency
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, style, align=PP_ALIGN.LEFT, color=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    current = dict(style)
    if color is not None:
        current["color"] = color
    set_run_style(run, current)
    return box


def add_bullets(slide, items, x, y, w, h, size=18, color=COLORS["mist"], gap=0.34):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.clear()
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"• {item}"
        set_run_style(run, font(size, False, color))
    return box


def add_panel(slide, x, y, w, h, fill=COLORS["navy"], transparency=0.18, radius=True, line_color=None, line_transparency=0.45):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    shape.line.color.rgb = line_color or COLORS["white"]
    shape.line.transparency = line_transparency
    shape.line.width = Pt(1.1)
    return shape


def add_chip(slide, text, x, y, w=None, color=COLORS["cyan"], fill=COLORS["white"], fill_transparency=0.88):
    w = w or max(1.15, min(3.2, 0.12 * len(text) + 0.75))
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = fill_transparency
    shape.line.color.rgb = color
    shape.line.transparency = 0.15
    shape.line.width = Pt(1.0)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_run_style(run, font(11, True, color))
    return shape


def add_picture_card(slide, image_path: Path, x, y, w, h, caption=None, accent=COLORS["cyan"]):
    shadow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.08), Inches(y + 0.1), Inches(w), Inches(h))
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = COLORS["navy"]
    shadow.fill.transparency = 0.42
    shadow.line.fill.background()

    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    frame.fill.solid()
    frame.fill.fore_color.rgb = COLORS["white"]
    frame.fill.transparency = 0.02
    frame.line.color.rgb = accent
    frame.line.transparency = 0.25
    frame.line.width = Pt(1.4)

    pad = 0.06
    slide.shapes.add_picture(str(image_path), Inches(x + pad), Inches(y + pad), width=Inches(w - 2 * pad), height=Inches(h - 2 * pad))
    if caption:
        cap = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.12), Inches(y + h - 0.48), Inches(min(w - 0.24, 2.85)), Inches(0.3))
        cap.fill.solid()
        cap.fill.fore_color.rgb = COLORS["navy"]
        cap.fill.transparency = 0.15
        cap.line.fill.background()
        tf = cap.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = caption
        set_run_style(run, font(10.5, True, COLORS["white"]))


def add_metric(slide, value, label, x, y, w=1.95):
    add_panel(slide, x, y, w, 1.06, fill=COLORS["white"], transparency=0.84, line_color=COLORS["cyan"], line_transparency=0.45)
    add_text(slide, value, x + 0.14, y + 0.16, w - 0.28, 0.33, font(20, True, COLORS["white"]), color=COLORS["white"])
    add_text(slide, label, x + 0.14, y + 0.56, w - 0.28, 0.28, font(10.5, False, COLORS["mist"]), color=COLORS["mist"])


def add_footer(slide, page_no, label="Project.Core™"):
    add_text(slide, label, 0.58, 7.05, 2.0, 0.2, font(9.5, True, COLORS["slate"]), color=COLORS["slate"])
    add_text(slide, str(page_no), 12.35, 7.03, 0.35, 0.2, font(9.5, True, COLORS["slate"]), align=PP_ALIGN.RIGHT, color=COLORS["slate"])


def add_title_block(slide, eyebrow, title, subtitle, left=0.72, top=0.54, width=6.2):
    add_text(slide, eyebrow, left, top, width, 0.24, font(10.5, True, COLORS["cyan"]), color=COLORS["cyan"])
    add_text(slide, title, left, top + 0.24, width, 1.25, font(26, True, COLORS["white"]), color=COLORS["white"])
    add_text(slide, subtitle, left, top + 1.35, width, 0.7, font(12.5, False, COLORS["mist"]), color=COLORS["mist"])


def build_presentation(backgrounds):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.author = "OpenAI Codex"
    prs.core_properties.title = "Project.Core™ — продуктовая презентация"
    prs.core_properties.subject = "Сайт и платформа предварительной бюджетной оценки систем безопасности"
    prs.core_properties.company = "Project.Core™"

    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_bg(slide, backgrounds["cover"])
    add_overlay(slide, COLORS["navy"], transparency=0.28)
    add_overlay(slide, COLORS["navy_soft"], transparency=0.22, x=0, y=0, w=7.0, h=7.5)
    add_chip(slide, "AI-платформа для пресейла и ТКП", 0.75, 0.54, w=2.55, color=COLORS["cyan"])
    add_text(slide, "Project.Core™", 0.75, 1.08, 4.8, 0.62, font(30, True, COLORS["white"]), color=COLORS["white"])
    add_text(
        slide,
        "Сайт и платформа предварительной бюджетной оценки\nсистем безопасности",
        0.75,
        1.62,
        5.8,
        1.2,
        font(24, True, COLORS["white"]),
        color=COLORS["white"],
    )
    add_text(
        slide,
        "Единый продукт для быстрого формирования бюджетной картины проекта, объяснения логики расчёта и подготовки коммерческого предложения.",
        0.75,
        2.95,
        5.45,
        0.9,
        font(13, False, COLORS["mist"]),
        color=COLORS["mist"],
    )
    chip_x = 0.75
    for text in ["6 подсистем", "5–10 минут", "85+ субъектов РФ", "AI-аудит цен", "Сайт + платформа"]:
        add_chip(slide, text, chip_x, 4.14, color=COLORS["orange"], fill=COLORS["white"], fill_transparency=0.9)
        chip_x += 1.42 if text == "5–10 минут" else 1.66
    add_picture_card(slide, SCREENSHOTS / "site-hero.png", 7.15, 0.78, 5.45, 5.92, caption="Главная страница продукта", accent=COLORS["orange"])
    add_text(slide, "Презентация по текущей версии разрабатываемого продукта", 0.75, 6.62, 5.0, 0.24, font(10.5, False, COLORS["slate"]), color=COLORS["slate"])

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_bg(slide, backgrounds["overview"])
    add_overlay(slide, COLORS["navy"], transparency=0.42)
    add_title_block(
        slide,
        "ПРОДУКТ В ОДНОМ КОНТУРЕ",
        "Project.Core™ объединяет публичный сайт и рабочую платформу",
        "Сайт объясняет продукт и заводит пользователя в демо-контур. Платформа собирает объект, считает бюджет и формирует выходные материалы.",
        width=6.5,
    )
    add_panel(slide, 0.74, 2.18, 5.15, 3.38, fill=COLORS["navy"], transparency=0.18)
    add_bullets(
        slide,
        [
            "Публичный сайт: позиционирование, демонстрация преимуществ, legal-контур, вход в демо.",
            "Рабочая платформа: объект, зоны, системы, проектирование, бюджет, логика расчёта, AI-риски.",
            "Результат: быстрый бюджет по объекту, который можно защищать на переговорах.",
            "Дополнительные выходы: ТКП, план проекта, спецификации и выгрузки.",
        ],
        1.0,
        2.5,
        4.6,
        2.6,
        size=16,
    )
    add_picture_card(slide, SCREENSHOTS / "site-hero.png", 6.35, 1.8, 6.05, 3.18, caption="Сайт продукта", accent=COLORS["cyan"])
    add_metric(slide, "6", "систем в одном расчёте", 6.35, 5.35)
    add_metric(slide, "5–10 мин", "среднее время оценки", 8.42, 5.35)
    add_metric(slide, "85+", "субъектов РФ в модели", 10.49, 5.35)
    add_footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_bg(slide, backgrounds["site"])
    add_overlay(slide, COLORS["navy"], transparency=0.48)
    add_title_block(
        slide,
        "ПУБЛИЧНЫЙ САЙТ",
        "Сайт показывает ценность продукта и переводит пользователя в платформу",
        "Маркетинговый контур сделан не как одностраничная витрина, а как часть продуктовой системы: от позиционирования до юридически прозрачного доступа.",
        width=5.9,
    )
    add_panel(slide, 0.72, 2.08, 4.95, 4.4, fill=COLORS["navy"], transparency=0.16)
    add_bullets(
        slide,
        [
            "Главный экран с УТП, метриками и входом в расчёт.",
            "Блоки «Почему это работает», сравнение с альтернативами и объяснение AI-движка.",
            "Страница «О системе» с подробным описанием модулей, алгоритмов и источников данных.",
            "Legal-контур: privacy, ПДн, cookies, пользовательское соглашение и disclaimer.",
        ],
        0.98,
        2.48,
        4.35,
        3.35,
        size=16,
    )
    add_picture_card(slide, SCREENSHOTS / "site-hero.png", 6.0, 1.72, 6.45, 2.55, caption="Главный экран", accent=COLORS["orange"])
    add_picture_card(slide, SCREENSHOTS / "site-comparison.png", 6.0, 4.45, 3.1, 2.05, caption="Позиционирование", accent=COLORS["cyan"])
    add_picture_card(slide, SCREENSHOTS / "site-ai-engine.png", 9.34, 4.45, 3.1, 2.05, caption="AI-блок", accent=COLORS["green"])
    add_footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_bg(slide, backgrounds["site_blocks"])
    add_overlay(slide, COLORS["navy"], transparency=0.5)
    add_title_block(
        slide,
        "БЛОКИ САЙТА",
        "На сайте уже видны ключевые сценарии демонстрации продукта",
        "Пользователь последовательно получает УТП, методологию расчёта, объяснение AI-контура и подробное описание системы.",
        width=8.2,
    )
    add_picture_card(slide, SCREENSHOTS / "site-hero.png", 0.72, 2.1, 3.95, 3.8, caption="УТП и вход в демо", accent=COLORS["orange"])
    add_picture_card(slide, SCREENSHOTS / "site-ai-engine.png", 4.87, 2.1, 3.95, 3.8, caption="Как работает AI-контур", accent=COLORS["cyan"])
    add_picture_card(slide, backgrounds["about_preview"], 9.02, 2.1, 3.6, 3.8, caption="Страница «О системе»", accent=COLORS["green"])
    add_text(slide, "Лендинг", 1.38, 6.1, 1.4, 0.2, font(11.5, True, COLORS["white"]), color=COLORS["white"])
    add_text(slide, "AI-движок", 5.61, 6.1, 1.6, 0.2, font(11.5, True, COLORS["white"]), color=COLORS["white"])
    add_text(slide, "Подробное описание", 9.62, 6.1, 2.0, 0.2, font(11.5, True, COLORS["white"]), color=COLORS["white"])
    add_footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_bg(slide, backgrounds["platform"])
    add_overlay(slide, COLORS["navy"], transparency=0.52)
    add_title_block(
        slide,
        "ПЛАТФОРМА: РАБОЧИЙ ИНТЕРФЕЙС",
        "Пользовательский контур платформы уже разбит на прикладные окна и шаги",
        "В интерфейсе есть объект, зонирование, состав систем, проектирование, бюджет, стоимость проекта, логика расчёта и AI-риски.",
        width=7.6,
    )
    add_picture_card(slide, SCREENSHOTS / "platform-object-view.png", 0.72, 2.0, 6.05, 4.75, caption="Окно «Объект»", accent=COLORS["cyan"])
    add_picture_card(slide, SCREENSHOTS / "platform-systems-view.png", 6.97, 2.0, 5.65, 4.75, caption="Окно «Системы»", accent=COLORS["orange"])
    add_chip(slide, "Зонирование", 0.86, 6.95, color=COLORS["cyan"])
    add_chip(slide, "AI-обследование", 2.25, 6.95, color=COLORS["cyan"])
    add_chip(slide, "PDF APS", 4.18, 6.95, color=COLORS["cyan"])
    add_chip(slide, "Вендоры и спецификация", 5.5, 6.95, color=COLORS["orange"], w=2.2)
    add_chip(slide, "Экспорт Excel / ТКП", 8.03, 6.95, color=COLORS["orange"], w=2.12)
    add_footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_bg(slide, backgrounds["platform2"])
    add_overlay(slide, COLORS["navy"], transparency=0.54)
    add_title_block(
        slide,
        "ФУНКЦИОНАЛ ПЛАТФОРМЫ",
        "Платформа считает, объясняет результат и заранее показывает риск-контур проекта",
        "На текущей версии уже присутствуют окна бюджета, стоимости проекта, логики расчёта и AI-рисков.",
        width=6.2,
    )
    add_panel(slide, 0.74, 2.18, 3.48, 4.55, fill=COLORS["navy"], transparency=0.18)
    add_bullets(
        slide,
        [
            "Управление коэффициентами и условиями расчёта бюджета.",
            "Агрегация стоимости по системам и по проекту в целом.",
            "Подробная объяснимость: что повлияло на сумму и трудозатраты.",
            "AI-риски проекта: критичные точки по спецификации, срокам и резервам.",
            "Выходные материалы: ТКП, план проекта, спецификации, презентации.",
        ],
        1.0,
        2.56,
        2.95,
        3.95,
        size=15.2,
    )
    add_picture_card(slide, SCREENSHOTS / "platform-budget-view.png", 4.52, 2.1, 3.78, 2.06, caption="Бюджет", accent=COLORS["cyan"])
    add_picture_card(slide, SCREENSHOTS / "platform-cost-view.png", 8.52, 2.1, 3.78, 2.06, caption="Стоимость проекта", accent=COLORS["orange"])
    add_picture_card(slide, SCREENSHOTS / "platform-logic-view.png", 4.52, 4.46, 3.78, 2.06, caption="Логика расчёта", accent=COLORS["green"])
    add_picture_card(slide, SCREENSHOTS / "platform-risks-view.png", 8.52, 4.46, 3.78, 2.06, caption="AI-риски", accent=COLORS["cyan"])
    add_footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_bg(slide, backgrounds["roadmap"])
    add_overlay(slide, COLORS["navy"], transparency=0.48)
    add_title_block(
        slide,
        "ПЛАНЫ ПО ДОРАБОТКАМ",
        "Следующий этап развития продукта после завершения основного кода",
        "Roadmap состоит из технического, правового и корпоративного контуров развития.",
        width=8.0,
    )
    roadmap_cards = [
        ("1. Завершение core-code", "Доведение основного функционала платформы и сайта до стабильной продуктовой версии."),
        ("2. Переход на отечественное ПО", "После завершения написания основного кода — замена и адаптация стека под целевой отечественный контур."),
        ("3. Правовая упаковка", "Регистрация товарного знака и самого продукта в соответствии с законодательством РФ."),
        ("4. Корпоративный secure-контур", "Выделенная доработка для ООО «СТК» (ПАО Сбер) с приведением в соответствие стандартам кибербезопасности РФ и ПАО Сбер."),
    ]
    x_positions = [0.76, 3.98, 7.2, 10.42]
    for idx, (title, text) in enumerate(roadmap_cards):
        add_panel(slide, x_positions[idx], 2.38, 2.55, 3.45, fill=COLORS["white"], transparency=0.83, line_color=COLORS["cyan"])
        add_text(slide, title, x_positions[idx] + 0.16, 2.62, 2.2, 0.62, font(14, True, COLORS["white"]), color=COLORS["white"])
        add_text(slide, text, x_positions[idx] + 0.16, 3.28, 2.16, 2.1, font(11.2, False, COLORS["mist"]), color=COLORS["mist"])
    add_footer(slide, 7)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_bg(slide, backgrounds["split"])
    add_overlay(slide, COLORS["navy"], transparency=0.54)
    add_title_block(
        slide,
        "ДВА ПРОДУКТА / ДВА КОНТУРА ПОСТАВКИ",
        "Развитие продукта разделяется на корпоративную ветку и внешний рынок",
        "Это позволяет одновременно решать задачи стратегического заказчика и строить масштабируемую коммерческую модель.",
        width=7.8,
    )
    add_panel(slide, 0.76, 2.1, 5.84, 4.86, fill=COLORS["navy"], transparency=0.14, line_color=COLORS["orange"])
    add_panel(slide, 6.73, 2.1, 5.84, 4.86, fill=COLORS["navy"], transparency=0.14, line_color=COLORS["cyan"])
    add_chip(slide, "Продукт 1", 0.98, 2.34, w=1.18, color=COLORS["orange"])
    add_text(slide, "Отдельная версия для ООО «СТК» (ПАО Сбер)", 0.98, 2.74, 4.95, 0.58, font(18, True, COLORS["white"]), color=COLORS["white"])
    add_bullets(
        slide,
        [
            "Работа под Генеральное соглашение с ПАО Сбер.",
            "Приведение решения в соответствие стандартам кибербезопасности РФ и ПАО Сбер.",
            "Отдельный контур внедрения, эксплуатации и сопровождения.",
            "Фокус на корпоративные требования, комплаенс и защищённость.",
        ],
        1.02,
        3.42,
        4.92,
        2.8,
        size=15,
    )
    add_chip(slide, "Продукт 2", 6.95, 2.34, w=1.18, color=COLORS["cyan"])
    add_text(slide, "Внешний рынок: доступ к платформе по подписке", 6.95, 2.74, 4.95, 0.58, font(18, True, COLORS["white"]), color=COLORS["white"])
    add_bullets(
        slide,
        [
            "Продажа доступа к платформе как коммерческого продукта.",
            "Тарифные планы с разной глубиной функционала и уровнем сервиса.",
            "Браузерный доступ без тяжёлого локального контура для клиента.",
            "Абонентская модель и масштабирование на несколько категорий заказчиков.",
        ],
        6.99,
        3.42,
        4.92,
        2.8,
        size=15,
    )
    add_footer(slide, 8)

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_bg(slide, backgrounds["tariffs"])
    add_overlay(slide, COLORS["navy"], transparency=0.5)
    add_title_block(
        slide,
        "ПРЕДВАРИТЕЛЬНАЯ МОДЕЛЬ ТАРИФОВ",
        "Внешний рынок: доступ к платформе по абонентской плате",
        "Ниже приведена рабочая схема тарификации, которую можно уточнить перед коммерческим запуском.",
        width=8.1,
    )
    tariffs = [
        ("Start", "от 49 000 ₽/мес", ["Базовый пресейл-контур", "Ограниченное число пользователей", "Стандартный экспорт результатов"], COLORS["cyan"]),
        ("Pro", "от 119 000 ₽/мес", ["Полный multi-system расчёт", "AI-обследование и AI-риски", "Расширенные выгрузки и план проекта"], COLORS["orange"]),
        ("Enterprise", "индивидуально", ["Корпоративная конфигурация", "Приоритетная поддержка и SLA", "Доработки и интеграции под заказчика"], COLORS["green"]),
    ]
    x_positions = [0.8, 4.42, 8.04]
    for idx, (name, price, features, accent) in enumerate(tariffs):
        add_panel(slide, x_positions[idx], 2.18, 3.18, 4.62, fill=COLORS["white"], transparency=0.82, line_color=accent, line_transparency=0.18)
        add_text(slide, name, x_positions[idx] + 0.18, 2.52, 2.8, 0.34, font(18, True, COLORS["white"]), color=COLORS["white"])
        add_text(slide, price, x_positions[idx] + 0.18, 2.98, 2.8, 0.38, font(17, True, accent), color=accent)
        add_bullets(slide, features, x_positions[idx] + 0.18, 3.62, 2.75, 2.08, size=13.2)
        add_chip(slide, "Абонентская модель", x_positions[idx] + 0.18, 6.18, w=1.82, color=accent)
    add_text(
        slide,
        "Для версии ООО «СТК» (ПАО Сбер) предполагается отдельная договорная модель поставки вне стандартной SaaS-тарификации.",
        0.86,
        6.96,
        11.4,
        0.22,
        font(10.5, False, COLORS["slate"]),
        color=COLORS["slate"],
    )
    add_footer(slide, 9)

    prs.save(str(OUTPUT))


if __name__ == "__main__":
    assets = build_assets()
    build_presentation(assets)
    print(OUTPUT)
