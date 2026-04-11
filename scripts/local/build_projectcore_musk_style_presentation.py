# -*- coding: utf-8 -*-
from pathlib import Path

from PIL import Image

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
ASSETS = ROOT / "presentation_assets"
SCREENSHOTS = ASSETS / "screenshots"
RENDERED = ASSETS / "rendered"
OUTPUT_DIR = ROOT / "presentations"
TMP_DIR = ROOT / "tmp_presentation"
OUTPUT = OUTPUT_DIR / "ProjectCore_Agent_Shell_Presentation.pptx"
OUTPUT_ALT = OUTPUT_DIR / "ProjectCore_Agent_Shell_Presentation_Musk_Style.pptx"
OUTPUT_V2 = OUTPUT_DIR / "ProjectCore_Agent_Shell_Presentation_Musk_Style_v2.pptx"
OUTPUT_V3 = OUTPUT_DIR / "ProjectCore_Agent_Shell_Presentation_Musk_Style_v3.pptx"

SW = 13.333
SH = 7.5

BLACK = RGBColor(4, 10, 18)
INK = RGBColor(11, 20, 32)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(176, 190, 210)
MUTED_SOFT = RGBColor(132, 151, 178)
CYAN = RGBColor(90, 214, 242)
BLUE = RGBColor(61, 132, 255)
GREEN = RGBColor(121, 230, 177)
ORANGE = RGBColor(255, 151, 89)
PANEL = RGBColor(15, 28, 44)


def set_font(run, size, bold=False, color=WHITE, name="Arial"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_shape(slide, x, y, w, h, color, transparency=0.0, rounded=False, line=None, line_transparency=0.0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.transparency = line_transparency
        shape.line.width = Pt(1.1)
    return shape


def add_bg(slide, image_path=None, color=BLACK):
    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), 0, 0, width=Inches(SW), height=Inches(SH))
    else:
        add_shape(slide, 0, 0, SW, SH, color)


def add_text(slide, text, x, y, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, bold, color)
    return box


def add_badge(slide, text, x, y, accent=CYAN, width=None):
    width = width or max(1.5, min(3.2, 0.105 * len(text) + 0.9))
    shape = add_shape(slide, x, y, width, 0.38, PANEL, transparency=0.18, rounded=True, line=accent, line_transparency=0.1)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font(run, 10.5, True, accent)
    return shape


def add_metric(slide, x, y, value, label, accent):
    add_shape(slide, x, y, 2.25, 1.1, PANEL, transparency=0.12, rounded=True, line=accent, line_transparency=0.12)
    add_text(slide, value, x + 0.18, y + 0.14, 1.7, 0.3, 24, True, WHITE)
    add_text(slide, label, x + 0.18, y + 0.62, 1.82, 0.25, 10.5, False, MUTED)


def add_full_overlay(slide, color=BLACK, transparency=0.35):
    add_shape(slide, 0, 0, SW, SH, color, transparency=transparency)


def add_image_contain(slide, image_path, x, y, w, h):
    with Image.open(image_path) as img:
        img_w, img_h = img.size

    box_ratio = w / h
    img_ratio = img_w / img_h

    if img_ratio > box_ratio:
        draw_w = w
        draw_h = w / img_ratio
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * img_ratio
        draw_x = x + (w - draw_w) / 2
        draw_y = y

    slide.shapes.add_picture(
        str(image_path),
        Inches(draw_x),
        Inches(draw_y),
        width=Inches(draw_w),
        height=Inches(draw_h),
    )


def prepare_cover_image(image_path, target_ratio, anchor="center"):
    TMP_DIR.mkdir(parents=True, exist_ok=True)
    image_path = Path(image_path)
    out_name = f"{image_path.stem}_cover_{str(target_ratio).replace('.', '_')}_{anchor}.png"
    out_path = TMP_DIR / out_name
    if out_path.exists():
        return out_path

    with Image.open(image_path) as img:
        img = img.convert("RGB")
        img_w, img_h = img.size
        img_ratio = img_w / img_h

        if img_ratio > target_ratio:
            new_w = int(img_h * target_ratio)
            if anchor == "left":
                left = 0
            elif anchor == "right":
                left = img_w - new_w
            else:
                left = (img_w - new_w) // 2
            crop = (left, 0, left + new_w, img_h)
        else:
            new_h = int(img_w / target_ratio)
            if anchor == "top":
                top = 0
            elif anchor == "bottom":
                top = img_h - new_h
            else:
                top = (img_h - new_h) // 2
            crop = (0, top, img_w, top + new_h)

        img.crop(crop).save(out_path, quality=95)

    return out_path


def add_device(slide, image_path, x, y, w, h, dark=True, screen_fill=RGBColor(244, 247, 250), fit="cover", anchor="center"):
    shadow_color = BLACK if dark else RGBColor(160, 171, 186)
    frame_color = PANEL if dark else WHITE
    line_color = RGBColor(45, 61, 83) if dark else RGBColor(220, 228, 236)
    add_shape(slide, x + 0.08, y + 0.12, w, h, shadow_color, transparency=0.74 if dark else 0.84, rounded=True)
    add_shape(slide, x, y, w, h, frame_color, transparency=0.0, rounded=True, line=line_color, line_transparency=0.15)
    inner_x = x + 0.08
    inner_y = y + 0.12
    inner_w = w - 0.16
    inner_h = h - 0.24
    add_shape(slide, inner_x, inner_y, inner_w, inner_h, screen_fill, rounded=True)
    if fit == "contain":
        add_image_contain(slide, image_path, inner_x + 0.04, inner_y + 0.04, inner_w - 0.08, inner_h - 0.08)
    else:
        prepared = prepare_cover_image(image_path, inner_w / inner_h, anchor=anchor)
        slide.shapes.add_picture(str(prepared), Inches(inner_x), Inches(inner_y), width=Inches(inner_w), height=Inches(inner_h))


def add_three_column_statement(slide, cols, y=4.55):
    xs = [0.86, 4.48, 8.1]
    accents = [CYAN, BLUE, GREEN]
    for i, (title, body) in enumerate(cols):
        add_shape(slide, xs[i], y, 3.18, 1.7, PANEL, transparency=0.14, rounded=True, line=accents[i], line_transparency=0.1)
        add_text(slide, title, xs[i] + 0.2, y + 0.18, 2.7, 0.3, 16, True, WHITE)
        add_text(slide, body, xs[i] + 0.2, y + 0.62, 2.74, 0.7, 11, False, MUTED)


def footer(slide, text="Spider0 | Product presentation"):
    add_text(slide, text, 0.72, 7.0, 4.0, 0.18, 9.5, False, RGBColor(214, 224, 236))


def build():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    bg_cover = RENDERED / "bg-cover.jpg"
    bg_platform = RENDERED / "bg-platform.jpg"
    bg_platform2 = RENDERED / "bg-platform2.jpg"
    bg_site = RENDERED / "bg-site.jpg"
    bg_split = RENDERED / "bg-split.jpg"
    bg_roadmap = RENDERED / "bg-roadmap.jpg"

    site_hero = SCREENSHOTS / "site-hero.png"
    site_ai = SCREENSHOTS / "site-ai-engine.png"
    site_about = RENDERED / "site-about-preview.jpg"
    platform_object = SCREENSHOTS / "platform-object-view.png"
    platform_logic = SCREENSHOTS / "platform-logic-view.png"
    platform_budget = SCREENSHOTS / "platform-budget-view.png"

    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_cover)
    add_full_overlay(slide, BLACK, 0.38)
    add_shape(slide, 0, 0, 6.55, SH, BLACK, transparency=0.08)
    add_badge(slide, "PRODUCT PRESENTATION", 0.82, 0.72, CYAN, 2.22)
    add_text(slide, "Один интерфейс.\nОдин агент.\nОдин рабочий контур.", 0.82, 1.32, 5.75, 2.05, 27, True, WHITE)
    add_text(
        slide,
        "ProjectCore Agent Shell объединяет AI-модели, файлы проекта, вложения, скриншоты и git-сценарии в одном сильном продукте для инженерной команды.",
        0.86,
        3.58,
        5.05,
        1.08,
        15,
        False,
        MUTED,
    )
    add_metric(slide, 0.86, 5.45, "3 модели", "GigaChat, Qwen CLI, Qwen API", CYAN)
    add_metric(slide, 3.28, 5.45, "1 поток", "контекст, запуск, результат", BLUE)
    footer(slide)

    # 2. Product thesis
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_platform2)
    add_full_overlay(slide, BLACK, 0.46)
    add_badge(slide, "WHY IT MATTERS", 0.82, 0.72, ORANGE, 1.72)
    add_text(slide, "Не ещё один AI-чат.", 0.82, 1.3, 5.8, 0.52, 27, True, WHITE)
    add_text(slide, "А AI-операционная система команды.", 0.82, 1.9, 7.8, 0.58, 27, True, WHITE)
    add_text(
        slide,
        "Смена провайдера, работа с проектом, вложения, промптинг и git больше не расползаются по окнам и утилитам. Всё собрано в одном управляемом интерфейсе.",
        0.86,
        3.1,
        5.8,
        1.06,
        15,
        False,
        MUTED,
    )
    add_three_column_statement(
        slide,
        [
            ("Быстрее старт", "От задачи до запуска агента без ручной подготовки окружения."),
            ("Меньше хаоса", "Файлы, скрины и доступы живут рядом, а не теряются по папкам."),
            ("Сильнее контроль", "Команда осознанно выбирает модель и git-контур под сценарий."),
        ],
        y=5.05,
    )
    footer(slide)

    # 3. Real product
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_site)
    add_full_overlay(slide, BLACK, 0.34)
    add_badge(slide, "REAL PRODUCT", 0.82, 0.7, CYAN, 1.52)
    add_text(slide, "Это уже выглядит как продукт.", 0.82, 1.25, 6.3, 0.5, 28, True, WHITE)
    add_text(slide, "Не как админка. Не как прототип. Не как набор скриптов.", 0.82, 1.9, 7.2, 0.34, 16, False, MUTED)
    add_device(slide, site_hero, 0.86, 2.34, 5.95, 4.4, dark=True, fit="cover", anchor="top")
    add_shape(slide, 7.08, 2.38, 5.34, 1.34, PANEL, transparency=0.14, rounded=True, line=CYAN, line_transparency=0.12)
    add_text(slide, "Публичная подача", 7.34, 2.64, 2.8, 0.24, 15, True, WHITE)
    add_text(slide, "С первого экрана понятно, в чём ценность продукта и почему ему можно доверять.", 7.34, 3.0, 4.78, 0.36, 10.8, False, MUTED)
    add_shape(slide, 7.08, 3.96, 5.34, 1.34, PANEL, transparency=0.14, rounded=True, line=BLUE, line_transparency=0.12)
    add_text(slide, "AI-позиционирование", 7.34, 4.22, 3.0, 0.24, 14, True, WHITE)
    add_text(slide, "Не магия, а понятный AI-движок, сценарий применения и инженерный смысл.", 7.34, 4.58, 4.8, 0.36, 10.8, False, MUTED)
    add_shape(slide, 7.08, 5.54, 5.34, 1.2, PANEL, transparency=0.14, rounded=True, line=GREEN, line_transparency=0.12)
    add_text(slide, "Готовность к продаже", 7.34, 5.78, 3.0, 0.24, 14, True, WHITE)
    add_text(slide, "Уже можно показывать клиенту, партнёру, инвестору и внутренней команде.", 7.34, 6.08, 4.82, 0.3, 10.8, False, MUTED)
    footer(slide)

    # 4. Functional stack
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_split)
    add_full_overlay(slide, BLACK, 0.42)
    add_badge(slide, "FUNCTIONAL STACK", 0.82, 0.72, CYAN, 1.9)
    add_text(slide, "Что делает продукт сильным", 0.82, 1.28, 5.4, 0.45, 27, True, WHITE)
    add_shape(slide, 0.82, 2.02, 5.35, 5.0, PANEL, transparency=0.12, rounded=True)
    add_text(slide, "Внутри одного интерфейса:", 1.06, 2.34, 2.9, 0.24, 18, True, WHITE)
    bullets = [
        "переключение между GigaChat, Qwen OAuth CLI и Qwen API",
        "привязка локальной рабочей папки и просмотр структуры проекта",
        "добавление файлов и скриншотов прямо в контекст задачи",
        "управление параметрами модели и доступами без ручной рутины",
        "переключение git remote между GitHub и GitVerse",
        "запуск агента и получение результата в том же окне",
    ]
    y = 2.9
    for item in bullets:
        add_shape(slide, 1.08, y + 0.08, 0.08, 0.08, CYAN, rounded=True)
        add_text(slide, item, 1.28, y, 4.5, 0.48, 14.5, False, WHITE)
        y += 0.7
    add_shape(slide, 7.25, 2.18, 5.15, 1.24, PANEL, transparency=0.14, rounded=True, line=CYAN, line_transparency=0.12)
    add_text(slide, "Модели", 7.52, 2.46, 1.2, 0.22, 17, True, WHITE)
    add_text(slide, "Команда выбирает нужный интеллект под задачу, бюджет и режим работы.", 7.52, 2.82, 4.42, 0.38, 11.2, False, MUTED)
    add_shape(slide, 7.25, 3.72, 5.15, 1.24, PANEL, transparency=0.14, rounded=True, line=BLUE, line_transparency=0.12)
    add_text(slide, "Контекст", 7.52, 3.9, 1.4, 0.22, 17, True, WHITE)
    add_text(slide, "Промпт, файлы и screenshots соединяются в один цельный сценарий работы.", 7.52, 4.36, 4.36, 0.38, 11.2, False, MUTED)
    add_shape(slide, 7.25, 5.26, 5.15, 1.24, PANEL, transparency=0.14, rounded=True, line=GREEN, line_transparency=0.12)
    add_text(slide, "Инженерный выход", 7.52, 5.34, 2.3, 0.22, 17, True, WHITE)
    add_text(slide, "Продукт доводит работу до результата, а не останавливается на красивом диалоге.", 7.52, 5.9, 4.4, 0.38, 11.2, False, MUTED)
    footer(slide)

    # 5. Screens of product
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_platform)
    add_full_overlay(slide, BLACK, 0.38)
    add_badge(slide, "PRODUCT SCREENS", 0.82, 0.72, ORANGE, 1.85)
    add_text(slide, "Сильная подача строится на реальном интерфейсе.", 0.82, 1.28, 7.5, 0.45, 26, True, WHITE)
    add_text(slide, "Поэтому в презентации нужны не абстрактные блоки, а экраны, которые подтверждают зрелость продукта.", 0.82, 1.88, 7.8, 0.34, 15.5, False, MUTED)
    add_text(slide, "Вместо трёх мелких превью здесь должен быть один читаемый экран. Поэтому главный интерфейс показываем крупно.", 0.82, 2.36, 9.6, 0.34, 13.2, False, MUTED)
    add_device(slide, platform_object, 0.82, 2.88, 9.2, 3.88, dark=True, fit="cover", anchor="top")
    add_shape(slide, 10.32, 2.92, 2.06, 1.04, PANEL, transparency=0.14, rounded=True, line=CYAN, line_transparency=0.12)
    add_text(slide, "Объект", 10.58, 3.18, 1.4, 0.2, 15, True, WHITE)
    add_text(slide, "Карточка проекта,\nпараметры и фото", 10.58, 3.5, 1.42, 0.34, 10.2, False, MUTED)
    add_shape(slide, 10.32, 4.16, 2.06, 1.04, PANEL, transparency=0.14, rounded=True, line=BLUE, line_transparency=0.12)
    add_text(slide, "Логика", 10.58, 4.42, 1.4, 0.2, 15, True, WHITE)
    add_text(slide, "Объяснение,\nкак формируется расчёт", 10.58, 4.74, 1.42, 0.34, 10.2, False, MUTED)
    add_shape(slide, 10.32, 5.4, 2.06, 1.04, PANEL, transparency=0.14, rounded=True, line=GREEN, line_transparency=0.12)
    add_text(slide, "Бюджет", 10.58, 5.66, 1.4, 0.2, 15, True, WHITE)
    add_text(slide, "Детализация цены\nи структуры стоимости", 10.58, 5.98, 1.42, 0.34, 10.2, False, MUTED)
    footer(slide)

    # 6. Why now
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_roadmap)
    add_full_overlay(slide, BLACK, 0.54)
    add_badge(slide, "WHY NOW", 0.82, 0.72, CYAN, 1.12)
    add_text(slide, "Рынок уже готов\nк такому продукту.", 0.82, 1.3, 5.8, 1.06, 29, True, WHITE)
    add_text(
        slide,
        "AI стал привычным. Но привычным не стал управляемый AI-процесс для команды. ProjectCore Agent Shell закрывает именно этот разрыв.",
        0.86,
        3.05,
        5.2,
        0.96,
        15,
        False,
        MUTED,
    )
    add_three_column_statement(
        slide,
        [
            ("AI есть у всех", "Но чаще всего он живёт как отдельный чат, а не как рабочая система."),
            ("Команды тонут в инструментах", "Скорость теряется на переключении контекста, а не на сложности задачи."),
            ("Нужны зрелые интерфейсы", "Побеждает не только модель, а то, как она встроена в реальную работу."),
        ],
        y=5.0,
    )
    footer(slide)

    # 7. Final slide
    slide = prs.slides.add_slide(blank)
    add_bg(slide, None, BLACK)
    add_shape(slide, 0, 0, SW, SH, RGBColor(6, 14, 24))
    add_badge(slide, "FINAL MESSAGE", 0.82, 0.72, GREEN, 1.65)
    add_text(slide, "ProjectCore Agent Shell\nвыглядит как продукт,\nа не как заготовка.", 0.82, 1.35, 6.4, 1.72, 29, True, WHITE)
    add_text(slide, "Именно так его и нужно показывать: крупно, смело, с реальными экранами и фотофонами, а не с мелкими карточками на пустом фоне.", 0.86, 3.72, 5.75, 0.96, 15, False, MUTED)
    add_shape(slide, 7.18, 1.26, 5.2, 5.2, PANEL, transparency=0.08, rounded=True, line=CYAN, line_transparency=0.12)
    add_device(slide, site_about, 7.46, 1.58, 4.64, 3.74, dark=True, fit="cover", anchor="top")
    add_text(slide, "Сайт тоже нужно показывать крупно. Один экран, одна мысль, один сильный визуальный акцент.", 7.54, 5.66, 4.2, 0.54, 11, False, MUTED)
    footer(slide)

    try:
        prs.save(str(OUTPUT))
        print(OUTPUT)
    except PermissionError:
        print(f"LOCKED: {OUTPUT}")

    try:
        prs.save(str(OUTPUT_ALT))
        print(OUTPUT_ALT)
    except PermissionError:
        print(f"LOCKED: {OUTPUT_ALT}")

    try:
        prs.save(str(OUTPUT_V2))
        print(OUTPUT_V2)
    except PermissionError:
        print(f"LOCKED: {OUTPUT_V2}")

    prs.save(str(OUTPUT_V3))
    print(OUTPUT_V3)


if __name__ == "__main__":
    build()
