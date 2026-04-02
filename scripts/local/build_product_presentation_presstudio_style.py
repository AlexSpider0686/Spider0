from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path.cwd()
ASSETS = ROOT / "presentation_assets"
SCREENSHOTS = ASSETS / "screenshots"
RENDERED = ASSETS / "rendered"
OUTPUT = ROOT / "ProjectCore_Product_Presentation_B2B_2026-03-30.pptx"

SW = 13.333
SH = 7.5

WHITE = RGBColor(255, 255, 255)
BG = RGBColor(244, 247, 251)
INK = RGBColor(25, 35, 52)
MUTED = RGBColor(104, 119, 143)
LINE = RGBColor(215, 223, 234)
BLUE = RGBColor(63, 104, 255)
CYAN = RGBColor(82, 200, 231)
ORANGE = RGBColor(255, 135, 78)
GREEN = RGBColor(86, 190, 137)
DARK = RGBColor(13, 23, 39)
DARK2 = RGBColor(21, 35, 56)
LIGHT_TEXT = RGBColor(233, 238, 246)


def apply_font(run, size, bold=False, color=INK, name="Aptos"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_shape(slide, x, y, w, h, color, transparency=0.0, rounded=False, line=None, line_transparency=0.0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.transparency = line_transparency
        shape.line.width = Pt(1)
    return shape


def add_bg(slide, image_path=None, color=BG):
    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), 0, 0, width=Inches(SW), height=Inches(SH))
    else:
        add_shape(slide, 0, 0, SW, SH, color)


def add_text(slide, text, x, y, w, h, size, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    apply_font(run, size, bold, color)
    return box


def add_bullets(slide, items, x, y, w, h, size=14, color=MUTED):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = f"• {item}"
        apply_font(run, size, False, color)
    return box


def add_chip(slide, text, x, y, accent=BLUE, width=None, dark=False):
    width = width or max(1.2, min(3.0, len(text) * 0.11 + 0.8))
    fill = WHITE if not dark else DARK2
    transparency = 0.03 if not dark else 0.18
    shape = add_shape(slide, x, y, width, 0.34, fill, transparency=transparency, rounded=True, line=accent, line_transparency=0.1)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    apply_font(run, 11, True, accent if not dark else WHITE)


def add_metric(slide, x, y, w, value, label, accent):
    add_shape(slide, x, y, w, 1.05, WHITE, rounded=True, line=LINE)
    add_shape(slide, x, y, 0.08, 1.05, accent)
    add_text(slide, value, x + 0.18, y + 0.12, w - 0.25, 0.26, 19, True, INK)
    add_text(slide, label, x + 0.18, y + 0.56, w - 0.25, 0.2, 10.5, False, MUTED)


def add_device(slide, image_path, x, y, w, h, dark=False):
    add_shape(slide, x + 0.08, y + 0.12, w, h, DARK if dark else RGBColor(181, 191, 205), transparency=0.76 if dark else 0.82, rounded=True)
    add_shape(slide, x, y, w, h, DARK2 if dark else WHITE, rounded=True, line=RGBColor(48, 60, 82) if dark else LINE)
    add_shape(slide, x + 0.12, y + 0.17, w - 0.24, h - 0.31, RGBColor(247, 249, 252), rounded=True)
    slide.shapes.add_picture(str(image_path), Inches(x + 0.14), Inches(y + 0.19), width=Inches(w - 0.28), height=Inches(h - 0.35))


def add_footer(slide, page, dark=False):
    color = LIGHT_TEXT if dark else MUTED
    add_text(slide, "Project.Core™", 0.56, 7.03, 1.8, 0.18, 9.5, True, color)
    add_text(slide, str(page), 12.2, 7.03, 0.35, 0.18, 9.5, True, color, align=PP_ALIGN.RIGHT)


def build():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[6]

    bg_cover = RENDERED / "bg-cover.jpg"
    bg_site = RENDERED / "bg-site.jpg"
    bg_platform = RENDERED / "bg-platform.jpg"
    bg_platform2 = RENDERED / "bg-platform2.jpg"
    bg_roadmap = RENDERED / "bg-roadmap.jpg"
    bg_split = RENDERED / "bg-split.jpg"
    bg_tariffs = RENDERED / "bg-tariffs.jpg"

    # 1. Титульный лист
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_cover)
    add_shape(slide, 0, 0, SW, SH, DARK, transparency=0.32)
    add_shape(slide, 0.68, 0.64, 5.92, 6.0, DARK2, transparency=0.16, rounded=True, line=WHITE, line_transparency=0.8)
    add_chip(slide, "Титульный лист", 0.95, 0.84, CYAN, width=1.32, dark=True)
    add_chip(slide, "B2B IT-продукт", 2.4, 0.84, ORANGE, width=1.45, dark=True)
    add_text(slide, "Project.Core™", 0.95, 1.34, 4.5, 0.42, 31, True, WHITE)
    add_text(slide, "Сайт и платформа\nпредварительной бюджетной оценки\nсистем безопасности", 0.95, 1.88, 4.9, 1.18, 25, True, WHITE)
    add_text(slide, "Продукт для быстрого формирования бюджетной картины проекта, демонстрации логики расчёта и подготовки коммерческого предложения для заказчика.", 0.95, 3.22, 4.95, 0.9, 13, False, LIGHT_TEXT)
    add_device(slide, SCREENSHOTS / "site-hero.png", 7.0, 0.76, 5.62, 5.98, dark=True)
    add_chip(slide, "6 систем", 0.96, 5.08, CYAN, width=1.15, dark=True)
    add_chip(slide, "5–10 минут", 2.25, 5.08, ORANGE, width=1.42, dark=True)
    add_chip(slide, "Сайт + платформа", 3.85, 5.08, GREEN, width=1.82, dark=True)
    add_footer(slide, 1, dark=True)

    # 2. Основное предложение
    slide = prs.slides.add_slide(blank)
    add_bg(slide, None, BG)
    add_text(slide, "Основное предложение", 0.68, 0.6, 4.0, 0.32, 28, True, INK)
    add_text(slide, "Project.Core™ — это продукт, который позволяет за 5–10 минут собрать предварительный бюджет по системам безопасности, увидеть логику формирования суммы и подготовить материал для переговоров с заказчиком.", 0.68, 1.02, 8.1, 0.62, 13, False, MUTED)
    add_metric(slide, 0.72, 1.9, 2.3, "5–10 мин", "время первичной оценки", BLUE)
    add_metric(slide, 3.18, 1.9, 2.3, "6", "подсистем в одном расчёте", CYAN)
    add_metric(slide, 5.64, 1.9, 2.3, "85+", "регионов РФ в модели", ORANGE)
    add_metric(slide, 8.1, 1.9, 2.3, "AI", "аудит цен и рисков", GREEN)
    add_shape(slide, 0.74, 3.18, 5.86, 3.56, WHITE, rounded=True, line=LINE)
    add_text(slide, "Что получает заказчик", 1.0, 3.48, 2.8, 0.22, 20, True, INK)
    add_bullets(slide, [
        "быстрый предварительный бюджет по объекту",
        "понятную расшифровку, из чего он собран",
        "снижение риска недооценки работ",
        "выходные материалы для ТКП и обсуждения проекта",
    ], 1.0, 3.94, 4.8, 1.9, size=14, color=MUTED)
    add_device(slide, SCREENSHOTS / "platform-object-view.png", 6.86, 3.22, 5.3, 3.42, dark=False)
    add_footer(slide, 2, dark=False)

    # 3. Проблематика
    slide = prs.slides.add_slide(blank)
    add_bg(slide, None, BG)
    add_text(slide, "Проблематика рынка", 0.68, 0.62, 4.0, 0.3, 28, True, INK)
    add_text(slide, "Большинство предварительных расчётов по системам безопасности по-прежнему делаются либо в тяжёлых сметных системах, либо в Excel-моделях, либо в разрозненных вендорских калькуляторах.", 0.68, 1.02, 8.0, 0.52, 12.8, False, MUTED)
    cards = [
        ("Тяжёлые сметные системы", "Сильны для ПСД и детальной сметы, но плохо подходят для быстрого пресейла.", BLUE),
        ("Excel-модели", "Гибкие, но зависят от автора, плохо версионируются и несут высокий риск ручной ошибки.", ORANGE),
        ("Вендорские калькуляторы", "Считают отдельные подсистемы, но не сводят весь проект в единую бюджетную картину.", CYAN),
    ]
    xs = [0.82, 4.42, 8.02]
    for i, (title, desc, accent) in enumerate(cards):
        add_shape(slide, xs[i], 2.08, 3.0, 2.58, WHITE, rounded=True, line=LINE)
        add_shape(slide, xs[i], 2.08, 0.08, 2.58, accent)
        add_text(slide, title, xs[i] + 0.18, 2.34, 2.55, 0.46, 17, True, INK)
        add_text(slide, desc, xs[i] + 0.18, 3.0, 2.52, 1.25, 12, False, MUTED)
    add_shape(slide, 1.58, 5.18, 10.2, 1.05, INK, transparency=0.0, rounded=True)
    add_text(slide, "Следствие для бизнеса: предварительный бюджет получается долго, непрозрачно и с риском недооценки.", 1.92, 5.52, 9.5, 0.22, 17, True, WHITE, align=PP_ALIGN.CENTER)
    add_footer(slide, 3, dark=False)

    # 4. Выгоды и преимущества
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_site)
    add_shape(slide, 0, 0, SW, SH, DARK, transparency=0.48)
    add_text(slide, "Не функции, а выгоды для бизнеса", 0.72, 0.58, 5.8, 0.34, 28, True, WHITE)
    add_text(slide, "По рекомендациям для продуктовых презентаций мы переводим функционал в понятную бизнес-ценность.", 0.72, 1.0, 6.7, 0.34, 12.5, False, LIGHT_TEXT)
    gains = [
        ("Единый расчёт по нескольким системам", "Вместо набора отдельных оценок заказчик видит общую бюджетную картину по объекту."),
        ("AI-аудит цен и трудозатрат", "Снижает риск занижения стоимости и делает предложение более защищаемым на переговорах."),
        ("Региональный и эксплуатационный контур", "Бюджет учитывает географию, статус здания и условия монтажа, а не считается в вакууме."),
        ("Explainability", "Пользователь видит происхождение суммы, а не только финальную цифру."),
    ]
    positions = [(0.84, 1.84), (6.92, 1.84), (0.84, 4.2), (6.92, 4.2)]
    accents = [BLUE, CYAN, ORANGE, GREEN]
    for i, ((title, desc), (x, y)) in enumerate(zip(gains, positions)):
        add_shape(slide, x, y, 5.58, 1.78, DARK2, transparency=0.16, rounded=True, line=accents[i], line_transparency=0.1)
        add_text(slide, title, x + 0.22, y + 0.24, 4.9, 0.24, 17, True, WHITE)
        add_text(slide, desc, x + 0.22, y + 0.72, 4.95, 0.7, 12.2, False, LIGHT_TEXT)
    add_footer(slide, 4, dark=True)

    # 5. Российский продукт / доверие
    slide = prs.slides.add_slide(blank)
    add_bg(slide, None, BG)
    add_text(slide, "Почему продукту можно доверять", 0.68, 0.62, 5.2, 0.32, 28, True, INK)
    add_text(slide, "Согласно выбранным референсам, доверие в B2B-презентации строится через востребованность, ясность применения и отдельный акцент на российский контур продукта.", 0.68, 1.02, 8.7, 0.46, 12.6, False, MUTED)
    add_shape(slide, 0.76, 1.9, 3.92, 4.8, WHITE, rounded=True, line=LINE)
    add_shape(slide, 4.88, 1.9, 3.92, 4.8, WHITE, rounded=True, line=LINE)
    add_shape(slide, 9.0, 1.9, 3.56, 4.8, WHITE, rounded=True, line=LINE)
    add_text(slide, "Востребованность", 1.0, 2.22, 2.2, 0.22, 20, True, INK)
    add_bullets(slide, [
        "быстрый пресейл для интеграторов",
        "подготовка ТКП и предварительной оценки",
        "ранняя защита бюджета перед заказчиком",
    ], 1.0, 2.76, 3.2, 1.5, size=14, color=MUTED)
    add_text(slide, "Сфера применения", 5.12, 2.22, 2.2, 0.22, 20, True, INK)
    add_bullets(slide, [
        "офисные и административные объекты",
        "общественные здания",
        "склады, производство, логистика",
        "объекты с несколькими подсистемами безопасности",
    ], 5.12, 2.76, 3.2, 1.7, size=14, color=MUTED)
    add_text(slide, "Российский контур", 9.24, 2.22, 2.2, 0.22, 20, True, INK)
    add_bullets(slide, [
        "учёт субъектов РФ",
        "отдельный roadmap перехода на отечественное ПО",
        "правовая регистрация продукта в РФ",
        "отдельная ветка под СТК / ПАО Сбер",
    ], 9.24, 2.76, 2.9, 1.95, size=13.4, color=MUTED)
    add_chip(slide, "Подходит для российского B2B-контекста", 4.24, 6.18, accent=ORANGE, width=2.85)
    add_footer(slide, 5, dark=False)

    # 6. Сайт
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_site)
    add_shape(slide, 0, 0, SW, SH, DARK, transparency=0.46)
    add_text(slide, "Сайт продукта", 0.72, 0.58, 3.4, 0.32, 28, True, WHITE)
    add_text(slide, "Публичная часть продукта объясняет ценность решения и подводит пользователя к демо и переговорам.", 0.72, 1.0, 6.1, 0.34, 12.5, False, LIGHT_TEXT)
    add_device(slide, SCREENSHOTS / "site-hero.png", 0.82, 1.72, 4.18, 4.96, dark=True)
    add_device(slide, SCREENSHOTS / "site-comparison.png", 5.16, 1.72, 3.34, 4.96, dark=True)
    add_device(slide, SCREENSHOTS / "site-ai-engine.png", 8.66, 1.72, 3.84, 4.96, dark=True)
    add_chip(slide, "Основное предложение", 1.52, 6.9, accent=CYAN, width=1.74, dark=True)
    add_chip(slide, "Преимущества и рынок", 5.76, 6.9, accent=ORANGE, width=1.92, dark=True)
    add_chip(slide, "AI-контур продукта", 9.76, 6.9, accent=GREEN, width=1.82, dark=True)
    add_footer(slide, 6, dark=True)

    # 7. Платформа
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_platform)
    add_shape(slide, 0, 0, SW, SH, DARK, transparency=0.5)
    add_text(slide, "Платформа: ключевые окна интерфейса", 0.72, 0.58, 5.7, 0.32, 28, True, WHITE)
    add_text(slide, "Рабочий интерфейс показывает не только поля ввода, но и структуру расчёта, логику, риски и готовые выходные материалы.", 0.72, 1.0, 7.1, 0.34, 12.5, False, LIGHT_TEXT)
    add_device(slide, SCREENSHOTS / "platform-object-view.png", 0.76, 1.82, 6.08, 4.96, dark=True)
    add_device(slide, SCREENSHOTS / "platform-systems-view.png", 6.98, 1.82, 5.62, 4.96, dark=True)
    add_chip(slide, "Объект, зоны, обследование", 1.22, 6.92, accent=CYAN, width=2.18, dark=True)
    add_chip(slide, "Системы, вендоры, спецификация", 7.42, 6.92, accent=ORANGE, width=2.42, dark=True)
    add_footer(slide, 7, dark=True)

    # 8. Функционал платформы
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_platform2)
    add_shape(slide, 0, 0, SW, SH, DARK, transparency=0.52)
    add_text(slide, "Что делает платформа", 0.72, 0.58, 4.4, 0.32, 28, True, WHITE)
    add_text(slide, "Функционал платформы лучше подавать как набор связанных окон и сценариев, а не как длинный список возможностей.", 0.72, 1.0, 7.2, 0.34, 12.5, False, LIGHT_TEXT)
    add_shape(slide, 0.8, 1.82, 3.24, 4.92, DARK2, transparency=0.18, rounded=True, line=WHITE, line_transparency=0.82)
    add_bullets(slide, [
        "объект и зонирование",
        "выбор состава систем",
        "вендоры и PDF APS",
        "бюджет и коэффициенты",
        "стоимость проекта",
        "логика расчёта",
        "AI-риски проекта",
        "экспорт ТКП и плана",
    ], 1.0, 2.12, 2.72, 4.1, size=14.1, color=LIGHT_TEXT)
    add_device(slide, SCREENSHOTS / "platform-budget-view.png", 4.34, 1.86, 3.98, 2.16, dark=True)
    add_device(slide, SCREENSHOTS / "platform-cost-view.png", 8.48, 1.86, 3.98, 2.16, dark=True)
    add_device(slide, SCREENSHOTS / "platform-logic-view.png", 4.34, 4.42, 3.98, 2.16, dark=True)
    add_device(slide, SCREENSHOTS / "platform-risks-view.png", 8.48, 4.42, 3.98, 2.16, dark=True)
    add_footer(slide, 8, dark=True)

    # 9. СТК/Сбер и внешний рынок
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_split)
    add_shape(slide, 0, 0, SW, SH, DARK, transparency=0.52)
    add_text(slide, "Две продуктовые линии", 0.72, 0.58, 4.2, 0.32, 28, True, WHITE)
    add_text(slide, "По логике коммерческого предложения важно сразу разделить корпоративный контур и рыночную модель продаж.", 0.72, 1.0, 7.0, 0.34, 12.5, False, LIGHT_TEXT)
    add_shape(slide, 0.78, 1.84, 5.84, 4.88, DARK2, transparency=0.18, rounded=True, line=ORANGE, line_transparency=0.12)
    add_shape(slide, 6.72, 1.84, 5.84, 4.88, DARK2, transparency=0.18, rounded=True, line=CYAN, line_transparency=0.12)
    add_chip(slide, "Продукт для ООО «СТК»", 1.0, 2.12, accent=ORANGE, width=1.95, dark=True)
    add_text(slide, "Отдельная разработка\nпод Генеральное соглашение\nс ПАО Сбер", 1.0, 2.52, 4.5, 1.0, 21, True, WHITE)
    add_bullets(slide, [
        "приведение в соответствие требованиям стандартов по кибербезопасности РФ и ПАО Сбер",
        "корпоративный контур внедрения и эксплуатации",
        "отдельный secure-roadmap продукта",
    ], 1.0, 3.9, 4.7, 1.9, size=14.2, color=LIGHT_TEXT)
    add_chip(slide, "Внешний рынок", 6.96, 2.12, accent=CYAN, width=1.38, dark=True)
    add_text(slide, "Продажа доступа\nк платформе\nпо подписке", 6.96, 2.52, 4.3, 1.0, 21, True, WHITE)
    add_bullets(slide, [
        "доступ к платформе как коммерческому сервису",
        "тарифные планы с разным функционалом",
        "абонентская плата как масштабируемая модель монетизации",
    ], 6.96, 3.9, 4.7, 1.9, size=14.2, color=LIGHT_TEXT)
    add_footer(slide, 9, dark=True)

    # 10. Тарифы и roadmap
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_tariffs)
    add_shape(slide, 0, 0, SW, SH, DARK, transparency=0.52)
    add_text(slide, "Тарифы и следующий этап развития", 0.72, 0.58, 6.4, 0.32, 28, True, WHITE)
    add_text(slide, "Внешний рынок предполагает SaaS-модель, а параллельно запланированы переход на отечественное ПО и правовая регистрация продукта в РФ.", 0.72, 1.0, 8.2, 0.38, 12.5, False, LIGHT_TEXT)
    tariffs = [
        ("Start", "от 49 000 ₽/мес", ["базовый пресейл", "стандартные выгрузки"], CYAN),
        ("Pro", "от 119 000 ₽/мес", ["полный расчёт", "AI-модули", "расширенный экспорт"], ORANGE),
        ("Enterprise", "индивидуально", ["корпоративная конфигурация", "поддержка и интеграции"], GREEN),
    ]
    xs = [0.88, 4.48, 8.08]
    for i, (name, price, features, accent) in enumerate(tariffs):
        add_shape(slide, xs[i], 1.84, 3.0, 3.88, DARK2, transparency=0.16, rounded=True, line=accent, line_transparency=0.1)
        add_text(slide, name, xs[i] + 0.2, 2.14, 1.6, 0.22, 20, True, WHITE)
        add_text(slide, price, xs[i] + 0.2, 2.56, 2.1, 0.22, 16, True, accent)
        add_bullets(slide, features, xs[i] + 0.2, 3.16, 2.4, 1.15, size=13.4, color=LIGHT_TEXT)
    add_shape(slide, 0.88, 6.0, 11.18, 0.88, WHITE, transparency=0.05, rounded=True, line=WHITE, line_transparency=0.75)
    add_text(slide, "После завершения процесса написания основного кода планируются: переход на отечественное ПО, регистрация товарного знака и регистрация самого продукта в соответствии с законодательством РФ.", 1.12, 6.26, 10.7, 0.3, 12.2, False, LIGHT_TEXT, align=PP_ALIGN.CENTER)
    add_footer(slide, 10, dark=True)

    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(path)
