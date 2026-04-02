from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path.cwd()
ASSETS = ROOT / "presentation_assets"
SCREENSHOTS = ASSETS / "screenshots"
OUTPUT = ROOT / "ProjectCore_Product_Presentation_Sber_Corporate_2026-03-30.pptx"

SW = 13.333
SH = 7.5

WHITE = RGBColor(255, 255, 255)
BG = RGBColor(246, 249, 247)
SURFACE = RGBColor(255, 255, 255)
TEXT = RGBColor(27, 38, 44)
MUTED = RGBColor(99, 113, 120)
LINE = RGBColor(214, 225, 219)
GREEN = RGBColor(33, 150, 83)
GREEN_DARK = RGBColor(25, 104, 58)
GREEN_SOFT = RGBColor(229, 243, 233)
NAVY = RGBColor(34, 52, 61)
GRAY_PANEL = RGBColor(241, 245, 243)


def font(run, size, bold=False, color=TEXT, name="Aptos"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def rect(slide, x, y, w, h, color, transparency=0.0, rounded=False, line=None, line_transparency=0.0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = line
        shape.line.transparency = line_transparency
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def text(slide, value, x, y, w, h, size, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    font(run, size, bold, color)
    return box


def bullets(slide, items, x, y, w, h, size=14, color=MUTED):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = f"• {item}"
        font(run, size, False, color)
    return box


def chip(slide, label, x, y, w=None, accent=GREEN):
    w = w or max(1.15, min(2.6, 0.11 * len(label) + 0.72))
    s = rect(slide, x, y, w, 0.34, GREEN_SOFT, rounded=True, line=accent, line_transparency=0.0)
    tf = s.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    font(run, 10.5, True, accent)


def metric(slide, x, y, w, value, label):
    rect(slide, x, y, w, 1.06, SURFACE, rounded=True, line=LINE)
    rect(slide, x, y, 0.08, 1.06, GREEN)
    text(slide, value, x + 0.18, y + 0.12, w - 0.25, 0.24, 20, True, TEXT)
    text(slide, label, x + 0.18, y + 0.58, w - 0.25, 0.18, 10.2, False, MUTED)


def device(slide, image_path, x, y, w, h):
    rect(slide, x + 0.06, y + 0.08, w, h, RGBColor(198, 210, 205), transparency=0.82, rounded=True)
    rect(slide, x, y, w, h, SURFACE, rounded=True, line=LINE)
    rect(slide, x + 0.1, y + 0.14, w - 0.2, h - 0.26, GRAY_PANEL, rounded=True)
    slide.shapes.add_picture(str(image_path), Inches(x + 0.12), Inches(y + 0.16), width=Inches(w - 0.24), height=Inches(h - 0.3))


def title_block(slide, eyebrow, title, subtitle):
    chip(slide, eyebrow, 0.7, 0.58, w=max(1.3, 0.11 * len(eyebrow) + 0.74))
    text(slide, title, 0.7, 0.98, 7.6, 0.44, 26, True, TEXT)
    text(slide, subtitle, 0.7, 1.38, 8.9, 0.46, 12.5, False, MUTED)


def footer(slide, page):
    text(slide, "Project.Core™", 0.56, 7.02, 1.6, 0.16, 9.5, True, MUTED)
    text(slide, str(page), 12.2, 7.02, 0.3, 0.16, 9.5, True, MUTED, align=PP_ALIGN.RIGHT)


def build():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[6]

    # 1 cover
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SW, SH, BG)
    rect(slide, 0.68, 0.72, 5.8, 5.98, SURFACE, rounded=True, line=LINE)
    rect(slide, 7.0, 0.72, 5.64, 5.98, SURFACE, rounded=True, line=LINE)
    chip(slide, "Титульный лист", 0.96, 0.96, w=1.34)
    text(slide, "Project.Core™", 0.96, 1.44, 4.4, 0.36, 28, True, TEXT)
    text(slide, "Сайт и платформа\nпредварительной бюджетной оценки\nсистем безопасности", 0.96, 1.94, 4.9, 1.05, 24, True, TEXT)
    text(slide, "Корпоративная продуктовая презентация в более строгой подаче: акцент на надёжность, структуру решения, применимость для enterprise-заказчика и отдельный контур для ООО «СТК» / ПАО Сбер.", 0.96, 3.26, 4.92, 1.0, 12.6, False, MUTED)
    metric(slide, 0.98, 4.82, 1.58, "6", "подсистем")
    metric(slide, 2.72, 4.82, 1.72, "5–10 мин", "предварительная оценка")
    metric(slide, 4.6, 4.82, 1.6, "85+", "субъектов РФ")
    device(slide, SCREENSHOTS / "platform-object-view.png", 7.16, 0.9, 5.32, 5.6)
    footer(slide, 1)

    # 2 executive summary
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SW, SH, BG)
    title_block(
        slide,
        "Executive Summary",
        "Project.Core™ — единый контур для презентации продукта и расчёта бюджета проекта",
        "Продукт объединяет публичный сайт и платформу расчёта, чтобы быстро подготовить предварительный бюджет, объяснить его происхождение и перевести обсуждение в рабочий диалог с заказчиком.",
    )
    metric(slide, 0.72, 1.92, 2.28, "Сайт", "демонстрация ценности и вход в демо")
    metric(slide, 3.18, 1.92, 2.28, "Платформа", "объект, системы, бюджет, AI")
    metric(slide, 5.64, 1.92, 2.28, "Enterprise", "корпоративный контур для СТК/Сбер")
    metric(slide, 8.1, 1.92, 2.28, "SaaS", "внешний рынок и тарифная модель")
    rect(slide, 0.72, 3.28, 5.9, 3.34, SURFACE, rounded=True, line=LINE)
    rect(slide, 6.74, 3.28, 5.88, 3.34, SURFACE, rounded=True, line=LINE)
    text(slide, "Основной результат", 0.98, 3.58, 2.4, 0.22, 19, True, TEXT)
    bullets(slide, [
        "ускорение подготовки предварительного бюджета",
        "снижение риска недооценки работ и трудозатрат",
        "повышение прозрачности расчёта для заказчика",
        "готовность к enterprise- и рыночному контуру поставки",
    ], 0.98, 4.06, 4.95, 1.8, size=14, color=MUTED)
    text(slide, "Текущая архитектура продукта", 7.0, 3.58, 3.2, 0.22, 19, True, TEXT)
    bullets(slide, [
        "публичный сайт с продуктовой и правовой подачей",
        "рабочая платформа с поэтапным процессом расчёта",
        "AI-модули аудита цен, логики и риск-контура",
        "экспорт ТКП, плана проекта и рабочих материалов",
    ], 7.0, 4.06, 4.9, 1.8, size=14, color=MUTED)
    footer(slide, 2)

    # 3 market problem
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SW, SH, BG)
    title_block(
        slide,
        "Проблематика",
        "Предварительный расчёт по системам безопасности остаётся медленным и фрагментированным",
        "На рынке преобладают тяжёлые сметные системы, Excel-модели и разрозненные калькуляторы отдельных вендоров.",
    )
    problems = [
        ("Тяжёлые сметные системы", "Подходят для детальной сметы, но избыточны для быстрого пресейла.", GREEN),
        ("Excel-модели", "Гибкие, но несут риск ошибки, зависят от автора и плохо масштабируются.", GREEN),
        ("Локальные калькуляторы", "Решают частную задачу, но не дают целостной картины проекта.", GREEN),
    ]
    xs = [0.86, 4.46, 8.06]
    for i, (head, body, accent) in enumerate(problems):
        rect(slide, xs[i], 2.18, 3.0, 2.44, SURFACE, rounded=True, line=LINE)
        rect(slide, xs[i], 2.18, 0.08, 2.44, accent)
        text(slide, head, xs[i] + 0.18, 2.46, 2.5, 0.4, 17, True, TEXT)
        text(slide, body, xs[i] + 0.18, 3.08, 2.48, 0.9, 12, False, MUTED)
    rect(slide, 0.88, 5.12, 11.7, 1.06, GREEN_DARK, rounded=True)
    text(slide, "Следствие: подготовка предварительного бюджета требует лишнего времени, а качество расчёта и защита суммы на переговорах зависят от ручного процесса.", 1.22, 5.42, 11.0, 0.22, 15.5, True, WHITE, align=PP_ALIGN.CENTER)
    footer(slide, 3)

    # 4 value
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SW, SH, BG)
    title_block(
        slide,
        "Ценность решения",
        "Project.Core™ переводит сложный расчёт в понятный и управляемый процесс",
        "В корпоративной подаче важны не отдельные функции, а прикладная ценность для заказчика и пресейл-команды.",
    )
    values = [
        ("Единая бюджетная картина", "Расчёт по нескольким подсистемам формируется в одном сценарии по объекту."),
        ("Снижение риска недооценки", "AI-аудит цен и трудозатрат защищает предложение от занижения."),
        ("Прозрачность суммы", "Пользователь видит, как сформировались объёмы, коэффициенты и итоговый бюджет."),
        ("Готовность к enterprise", "Продукт уже укладывается в логику корпоративной и рыночной поставки."),
    ]
    coords = [(0.82, 2.0), (6.84, 2.0), (0.82, 4.28), (6.84, 4.28)]
    for (head, body), (x, y) in zip(values, coords):
        rect(slide, x, y, 5.66, 1.76, SURFACE, rounded=True, line=LINE)
        rect(slide, x, y, 0.08, 1.76, GREEN)
        text(slide, head, x + 0.2, y + 0.22, 4.7, 0.22, 17, True, TEXT)
        text(slide, body, x + 0.2, y + 0.72, 4.9, 0.56, 12.1, False, MUTED)
    footer(slide, 4)

    # 5 site
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SW, SH, BG)
    title_block(
        slide,
        "Публичный сайт",
        "Сайт формирует доверие к продукту и открывает вход в рабочий сценарий",
        "Публичная часть продукта показывает основное предложение, логику решения и правовой контур использования.",
    )
    device(slide, SCREENSHOTS / "site-hero.png", 0.82, 1.96, 4.18, 4.8)
    device(slide, SCREENSHOTS / "site-comparison.png", 5.16, 1.96, 3.26, 4.8)
    device(slide, SCREENSHOTS / "site-ai-engine.png", 8.58, 1.96, 3.76, 4.8)
    chip(slide, "Главный экран и УТП", 1.54, 6.92, w=1.72)
    chip(slide, "Позиционирование", 5.86, 6.92, w=1.5)
    chip(slide, "AI-логика и доверие", 9.46, 6.92, w=1.94)
    footer(slide, 5)

    # 6 platform
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SW, SH, BG)
    title_block(
        slide,
        "Платформа",
        "Платформа расчёта выстроена как поэтапный рабочий контур",
        "Пользователь проходит от параметров объекта и систем до бюджетной картины, логики расчёта и риск-контуров проекта.",
    )
    device(slide, SCREENSHOTS / "platform-object-view.png", 0.78, 1.92, 6.0, 4.86)
    device(slide, SCREENSHOTS / "platform-systems-view.png", 6.96, 1.92, 5.58, 4.86)
    chip(slide, "Объект, зоны, обследование", 1.18, 6.92, w=2.16)
    chip(slide, "Системы, вендоры, спецификация", 7.4, 6.92, w=2.48)
    footer(slide, 6)

    # 7 functionality
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SW, SH, BG)
    title_block(
        slide,
        "Ключевой функционал платформы",
        "Платформа сочетает расчёт, explainability и экспортные сценарии",
        "Слайд показывает не полный список опций, а основные функциональные контуры, важные для B2B-поставки.",
    )
    rect(slide, 0.82, 1.98, 3.12, 4.74, SURFACE, rounded=True, line=LINE)
    bullets(slide, [
        "объект, адрес, площадь, этажность, регион",
        "распределение зон",
        "реестр систем и вендоров",
        "PDF APS и спецификации",
        "бюджет и коэффициенты",
        "стоимость проекта",
        "логика расчёта",
        "AI-риски проекта",
        "экспорт ТКП, плана, Excel",
    ], 1.02, 2.26, 2.56, 4.0, size=13.8, color=MUTED)
    device(slide, SCREENSHOTS / "platform-budget-view.png", 4.28, 1.98, 4.0, 2.14)
    device(slide, SCREENSHOTS / "platform-cost-view.png", 8.48, 1.98, 4.0, 2.14)
    device(slide, SCREENSHOTS / "platform-logic-view.png", 4.28, 4.5, 4.0, 2.14)
    device(slide, SCREENSHOTS / "platform-risks-view.png", 8.48, 4.5, 4.0, 2.14)
    footer(slide, 7)

    # 8 roadmap + domestic
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SW, SH, BG)
    title_block(
        slide,
        "План развития продукта",
        "Следующий этап — enterprise-готовность и нормативная упаковка",
        "После завершения основного кода развитие продукта включает технологический, правовой и корпоративный контуры.",
    )
    items = [
        ("01", "Завершение основного контура разработки", "Стабилизация сайта и платформы как единого продукта."),
        ("02", "Переход на отечественное ПО", "Запланирован после завершения процесса написания основного кода."),
        ("03", "Регистрация товарного знака и продукта", "Правовое оформление в соответствии с законодательством РФ."),
        ("04", "Корпоративный secure-контур", "Адаптация под требования стандартов кибербезопасности РФ и ПАО Сбер."),
    ]
    y = 1.94
    for n, head, body in items:
        rect(slide, 0.86, y, 11.56, 1.0, SURFACE, rounded=True, line=LINE)
        text(slide, n, 1.08, y + 0.2, 0.38, 0.2, 17, True, GREEN_DARK)
        text(slide, head, 1.72, y + 0.16, 4.4, 0.22, 16.5, True, TEXT)
        text(slide, body, 6.08, y + 0.18, 5.8, 0.22, 11.8, False, MUTED)
        y += 1.18
    footer(slide, 8)

    # 9 dual contour
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SW, SH, BG)
    title_block(
        slide,
        "Две продуктовые линии",
        "Корпоративный контур для СТК/ПАО Сбер и отдельный внешний рынок",
        "Продукт целесообразно развивать в двух самостоятельных направлениях поставки и эксплуатации.",
    )
    rect(slide, 0.82, 2.0, 5.72, 4.76, SURFACE, rounded=True, line=LINE)
    rect(slide, 6.8, 2.0, 5.72, 4.76, SURFACE, rounded=True, line=LINE)
    chip(slide, "Корпоративный контур", 1.04, 2.24, w=1.92)
    text(slide, "Отдельный продукт для ООО «СТК»\nпод Генеральное соглашение с ПАО Сбер", 1.04, 2.66, 4.72, 0.8, 20, True, TEXT)
    bullets(slide, [
        "соответствие требованиям стандартов по кибербезопасности РФ",
        "приведение в соответствие стандартам ПАО Сбер",
        "отдельный сценарий внедрения, эксплуатации и сопровождения",
    ], 1.04, 3.82, 4.92, 1.75, size=14, color=MUTED)
    chip(slide, "Внешний рынок", 7.02, 2.24, w=1.36)
    text(slide, "Коммерческий продукт\nс доступом к платформе\nпо подписке", 7.02, 2.66, 4.52, 0.8, 20, True, TEXT)
    bullets(slide, [
        "продажа доступа к платформе как сервису",
        "тарифные планы с разным уровнем функционала",
        "абонентская плата как масштабируемая модель монетизации",
    ], 7.02, 3.82, 4.92, 1.75, size=14, color=MUTED)
    footer(slide, 9)

    # 10 tariffs
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, SW, SH, BG)
    title_block(
        slide,
        "Тарифная модель внешнего рынка",
        "Разный объём функционала — разный уровень доступа и стоимости",
        "Внешний рынок предполагает подписочную модель доступа к платформе, тогда как корпоративный контур поставляется отдельно.",
    )
    plans = [
        ("Start", "от 49 000 ₽/мес", ["базовый пресейл", "стандартные выгрузки", "ограниченный доступ"]),
        ("Pro", "от 119 000 ₽/мес", ["полный расчёт", "AI-модули", "расширенные материалы"]),
        ("Enterprise", "индивидуально", ["корпоративная конфигурация", "поддержка и SLA", "интеграции"]),
    ]
    xs = [0.9, 4.5, 8.1]
    for i, (head, price, feats) in enumerate(plans):
        rect(slide, xs[i], 1.98, 3.0, 4.64, SURFACE, rounded=True, line=LINE)
        text(slide, head, xs[i] + 0.2, 2.28, 1.6, 0.22, 19, True, TEXT)
        text(slide, price, xs[i] + 0.2, 2.72, 2.0, 0.22, 15.5, True, GREEN_DARK)
        bullets(slide, feats, xs[i] + 0.2, 3.3, 2.4, 1.5, size=13.6, color=MUTED)
        chip(slide, "Абонентская плата", xs[i] + 0.2, 5.92, w=1.7)
    text(slide, "Для версии ООО «СТК» / ПАО Сбер предполагается отдельная договорная модель поставки вне стандартной тарифной линейки.", 0.92, 6.96, 11.3, 0.2, 10.5, False, MUTED)
    footer(slide, 10)

    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    result = build()
    print(result)
