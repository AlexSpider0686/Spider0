from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path.cwd()
ASSETS = ROOT / "presentation_assets"
SCREENSHOTS = ASSETS / "screenshots"
RENDERED = ASSETS / "rendered"
OUTPUT = ROOT / "ProjectCore_Product_Presentation_Premium_2026-03-30.pptx"

SW = 13.333
SH = 7.5

WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(9, 15, 24)
NAVY = RGBColor(10, 18, 31)
NAVY_2 = RGBColor(17, 31, 51)
TEXT = RGBColor(232, 238, 245)
TEXT_DARK = RGBColor(23, 33, 51)
MUTED = RGBColor(167, 181, 200)
MUTED_DARK = RGBColor(99, 115, 137)
LINE = RGBColor(210, 221, 233)
CYAN = RGBColor(90, 214, 242)
ORANGE = RGBColor(255, 136, 74)
GREEN = RGBColor(111, 214, 157)
BLUE = RGBColor(75, 122, 255)
BG = RGBColor(245, 248, 252)


def set_font(run, size, bold=False, color=WHITE, name="Aptos"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_rect(slide, x, y, w, h, color, transparency=0.0, rounded=False, line=None, line_transparency=0.0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = line
        shape.line.transparency = line_transparency
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_bg(slide, image=None, color=BG):
    if image and Path(image).exists():
        slide.shapes.add_picture(str(image), 0, 0, width=Inches(SW), height=Inches(SH))
    else:
        add_rect(slide, 0, 0, SW, SH, color)


def add_text(slide, text, x, y, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, bold, color)
    return box


def add_bullets(slide, items, x, y, w, h, size=15, color=TEXT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = f"• {item}"
        set_font(run, size, False, color)
    return box


def add_chip(slide, text, x, y, accent=CYAN, width=None, light=False):
    width = width or max(1.15, min(3.0, 0.11 * len(text) + 0.8))
    fill = WHITE if light else NAVY
    transparency = 0.08 if light else 0.28
    line_color = accent
    shape = add_rect(slide, x, y, width, 0.34, fill, transparency=transparency, rounded=True, line=line_color, line_transparency=0.15)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font(run, 11, True, accent if light else WHITE)


def add_metric(slide, x, y, w, value, label, accent):
    add_rect(slide, x, y, w, 1.1, WHITE, transparency=0.03, rounded=True, line=LINE)
    add_rect(slide, x, y, 0.08, 1.1, accent, rounded=False)
    add_text(slide, value, x + 0.18, y + 0.13, w - 0.24, 0.3, 21, True, TEXT_DARK)
    add_text(slide, label, x + 0.18, y + 0.62, w - 0.24, 0.22, 10.5, False, MUTED_DARK)


def add_device(slide, image_path, x, y, w, h, dark=True):
    add_rect(slide, x + 0.08, y + 0.12, w, h, BLACK if dark else RGBColor(170, 181, 197), transparency=0.72 if dark else 0.82, rounded=True)
    add_rect(slide, x, y, w, h, NAVY if dark else WHITE, rounded=True, line=RGBColor(50, 62, 83) if dark else LINE)
    add_rect(slide, x + 0.12, y + 0.18, w - 0.24, h - 0.33, RGBColor(245, 248, 252), rounded=True)
    slide.shapes.add_picture(str(image_path), Inches(x + 0.14), Inches(y + 0.2), width=Inches(w - 0.28), height=Inches(h - 0.37))


def footer(slide, page):
    add_text(slide, "Project.Core™", 0.56, 7.02, 1.8, 0.18, 9.5, True, MUTED if page in [1, 3, 4, 6, 7, 8, 9, 10] else MUTED_DARK)
    add_text(slide, str(page), 12.24, 7.02, 0.35, 0.18, 9.5, True, MUTED if page in [1, 3, 4, 6, 7, 8, 9, 10] else MUTED_DARK, align=PP_ALIGN.RIGHT)


def build():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[6]

    bg_cover = RENDERED / "bg-cover.jpg"
    bg_site = RENDERED / "bg-site.jpg"
    bg_platform = RENDERED / "bg-platform.jpg"
    bg_platform2 = RENDERED / "bg-platform2.jpg"
    bg_roadmap = RENDERED / "bg-roadmap.jpg"
    bg_split = RENDERED / "bg-split.jpg"
    bg_tariffs = RENDERED / "bg-tariffs.jpg"

    # 1 титульный лист
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_cover)
    add_rect(slide, 0, 0, SW, SH, NAVY, transparency=0.34)
    add_rect(slide, 0.62, 0.58, 5.9, 6.08, NAVY, transparency=0.15, rounded=True, line=WHITE, line_transparency=0.78)
    add_chip(slide, "Титульный лист", 0.9, 0.82, accent=CYAN, width=1.38)
    add_chip(slide, "Сайт + платформа", 2.42, 0.82, accent=ORANGE, width=1.82)
    add_text(slide, "Project.Core™", 0.9, 1.35, 4.5, 0.45, 31, True, WHITE)
    add_text(slide, "Презентация\nразрабатываемого продукта", 0.9, 1.9, 4.7, 0.95, 28, True, WHITE)
    add_text(slide, "Система предварительной бюджетной оценки систем безопасности с современным публичным сайтом и рабочей платформой для расчёта, объяснения результата и подготовки коммерческого предложения.", 0.9, 3.03, 5.0, 1.02, 13.2, False, TEXT)
    add_metric(slide, 0.92, 4.72, 1.72, "6", "подсистем", CYAN)
    add_metric(slide, 2.78, 4.72, 1.72, "5–10 мин", "средняя оценка", ORANGE)
    add_metric(slide, 4.64, 4.72, 1.72, "85+", "субъектов РФ", GREEN)
    add_device(slide, SCREENSHOTS / "site-hero.png", 6.98, 0.74, 5.68, 6.05, dark=True)
    footer(slide, 1)

    # 2 обзор
    slide = prs.slides.add_slide(blank)
    add_bg(slide, None, BG)
    add_text(slide, "Что представляет собой продукт", 0.68, 0.6, 5.8, 0.36, 28, True, TEXT_DARK)
    add_text(slide, "Project.Core™ — это не только интерфейс расчёта, а полноценный продуктовый контур: презентационный сайт, рабочая платформа и дорожная карта вывода на корпоративный и внешний рынок.", 0.68, 1.0, 8.2, 0.52, 12.6, False, MUTED_DARK)
    add_metric(slide, 0.7, 1.78, 2.35, "Сайт", "презентация ценности и вход в демо", BLUE)
    add_metric(slide, 3.22, 1.78, 2.35, "Платформа", "объект, системы, бюджет, AI", CYAN)
    add_metric(slide, 5.74, 1.78, 2.35, "СТК / Сбер", "отдельный корпоративный контур", ORANGE)
    add_metric(slide, 8.26, 1.78, 2.35, "SaaS", "продажа доступа на рынок", GREEN)
    add_rect(slide, 0.72, 3.12, 5.9, 3.84, WHITE, rounded=True, line=LINE)
    add_rect(slide, 6.72, 3.12, 5.9, 3.84, WHITE, rounded=True, line=LINE)
    add_text(slide, "Публичный контур", 1.0, 3.42, 2.2, 0.25, 20, True, TEXT_DARK)
    add_text(slide, "Сайт показывает, почему продукт нужен рынку, как работает методология и за счёт чего обеспечивается доверие к расчёту.", 1.0, 3.78, 4.8, 0.52, 12, False, MUTED_DARK)
    add_bullets(slide, [
        "hero-экран и позиционирование",
        "объяснение AI-движка",
        "страница «О системе»",
        "legal-контур"
    ], 1.0, 4.52, 2.8, 1.6, size=13.2, color=MUTED_DARK)
    add_text(slide, "Рабочий контур", 7.0, 3.42, 2.2, 0.25, 20, True, TEXT_DARK)
    add_text(slide, "Платформа собирает объект, состав систем, бюджеты, AI-риски и формирует материалы для дальнейшей работы с заказчиком.", 7.0, 3.78, 4.8, 0.52, 12, False, MUTED_DARK)
    add_bullets(slide, [
        "объект и зонирование",
        "системы и спецификация",
        "бюджет и explainability",
        "экспорт ТКП и плана"
    ], 7.0, 4.52, 2.8, 1.6, size=13.2, color=MUTED_DARK)
    add_device(slide, SCREENSHOTS / "site-hero.png", 3.6, 4.22, 2.55, 2.26, dark=False)
    add_device(slide, SCREENSHOTS / "platform-object-view.png", 9.0, 4.22, 2.75, 2.26, dark=False)
    footer(slide, 2)

    # 3 сайт
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_site)
    add_rect(slide, 0, 0, SW, SH, NAVY, transparency=0.46)
    add_text(slide, "Публичный сайт", 0.72, 0.58, 3.5, 0.34, 28, True, WHITE)
    add_text(slide, "Сайт должен выглядеть как современный software-product showcase: быстро объяснять ценность, подтверждать доверие и переводить пользователя в платформу.", 0.72, 1.0, 6.2, 0.54, 12.5, False, TEXT)
    add_bullets(slide, [
        "УТП и CTA на первом экране.",
        "Сравнение с альтернативами рынка.",
        "Пояснение AI-аудита цен и трудозатрат.",
        "Подробное описание системы и правовой контур."
    ], 0.88, 1.82, 4.7, 1.9, size=15)
    add_device(slide, SCREENSHOTS / "site-hero.png", 6.15, 0.82, 6.28, 2.62, dark=True)
    add_device(slide, SCREENSHOTS / "site-comparison.png", 6.15, 3.72, 3.06, 2.9, dark=True)
    add_device(slide, SCREENSHOTS / "site-ai-engine.png", 9.46, 3.72, 2.97, 2.9, dark=True)
    footer(slide, 3)

    # 4 блоки сайта
    slide = prs.slides.add_slide(blank)
    add_bg(slide, None, BG)
    add_text(slide, "Что уже есть на сайте", 0.68, 0.62, 4.0, 0.34, 28, True, TEXT_DARK)
    add_text(slide, "Слайды ниже показывают не абстрактные макеты, а реальные блоки текущего проекта: лендинг, AI-контур и подробную страницу о платформе.", 0.68, 1.04, 7.4, 0.44, 12.4, False, MUTED_DARK)
    add_device(slide, SCREENSHOTS / "site-hero.png", 0.78, 1.82, 4.05, 4.88, dark=False)
    add_device(slide, SCREENSHOTS / "site-ai-engine.png", 4.97, 1.82, 4.05, 4.88, dark=False)
    add_device(slide, RENDERED / "site-about-preview.jpg", 9.16, 1.82, 3.42, 4.88, dark=False)
    add_chip(slide, "Лендинг и позиционирование", 1.34, 6.9, accent=BLUE, width=2.15, light=True)
    add_chip(slide, "AI-блок и методология", 5.58, 6.9, accent=CYAN, width=1.95, light=True)
    add_chip(slide, "Страница «О системе»", 9.78, 6.9, accent=GREEN, width=1.95, light=True)
    footer(slide, 4)

    # 5 платформа
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_platform)
    add_rect(slide, 0, 0, SW, SH, NAVY, transparency=0.5)
    add_text(slide, "Платформа: рабочие окна интерфейса", 0.72, 0.58, 5.8, 0.34, 28, True, WHITE)
    add_text(slide, "Внутри платформы уже выстроен последовательный рабочий сценарий: от описания объекта до результата, который можно использовать в коммерческой и инженерной работе.", 0.72, 1.0, 7.2, 0.54, 12.5, False, TEXT)
    add_device(slide, SCREENSHOTS / "platform-object-view.png", 0.72, 1.86, 6.12, 4.92, dark=True)
    add_device(slide, SCREENSHOTS / "platform-systems-view.png", 6.98, 1.86, 5.62, 4.92, dark=True)
    add_chip(slide, "Окно «Объект»", 1.1, 6.92, accent=CYAN, width=1.55)
    add_chip(slide, "Окно «Системы»", 7.4, 6.92, accent=ORANGE, width=1.62)
    footer(slide, 5)

    # 6 функционал платформы
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_platform2)
    add_rect(slide, 0, 0, SW, SH, NAVY, transparency=0.52)
    add_text(slide, "Функционал платформы", 0.72, 0.58, 4.5, 0.34, 28, True, WHITE)
    add_text(slide, "Платформа не только считает сумму, но и показывает логику, риски и выходные артефакты проекта.", 0.72, 1.0, 5.9, 0.38, 12.5, False, TEXT)
    add_rect(slide, 0.78, 1.8, 3.3, 4.95, NAVY_2, transparency=0.18, rounded=True, line=WHITE, line_transparency=0.75)
    add_bullets(slide, [
        "объект, адрес, площадь, этажность, регион",
        "зоны и шаблоны распределения",
        "системы, вендоры, PDF APS",
        "бюджет и коэффициенты",
        "стоимость проекта",
        "логика расчёта",
        "AI-риски проекта",
        "экспорт ТКП, Excel, план проекта"
    ], 1.0, 2.18, 2.75, 4.2, size=14.2)
    add_device(slide, SCREENSHOTS / "platform-budget-view.png", 4.42, 1.84, 3.95, 2.18, dark=True)
    add_device(slide, SCREENSHOTS / "platform-cost-view.png", 8.56, 1.84, 3.95, 2.18, dark=True)
    add_device(slide, SCREENSHOTS / "platform-logic-view.png", 4.42, 4.38, 3.95, 2.18, dark=True)
    add_device(slide, SCREENSHOTS / "platform-risks-view.png", 8.56, 4.38, 3.95, 2.18, dark=True)
    footer(slide, 6)

    # 7 доработки
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_roadmap)
    add_rect(slide, 0, 0, SW, SH, NAVY, transparency=0.48)
    add_text(slide, "Планы по доработкам", 0.72, 0.58, 4.5, 0.34, 28, True, WHITE)
    add_text(slide, "Следующий этап после завершения процесса написания основного кода.", 0.72, 1.0, 5.5, 0.32, 12.5, False, TEXT)
    roadmap = [
        ("01", "Завершение core-функционала", "Стабилизация основного контура сайта и платформы."),
        ("02", "Переход на отечественное ПО", "Запланирован после завершения основной разработки."),
        ("03", "Регистрация ТЗ и продукта", "Приведение в соответствие требованиям законодательства РФ."),
        ("04", "Корпоративная secure-ветка", "Отдельная доработка для ООО «СТК» (ПАО Сбер)."),
    ]
    y = 1.82
    accents = [CYAN, ORANGE, GREEN, BLUE]
    for i, item in enumerate(roadmap):
        n, title, desc = item
        add_rect(slide, 0.86, y, 11.58, 1.08, NAVY_2, transparency=0.18, rounded=True, line=accents[i], line_transparency=0.15)
        add_text(slide, n, 1.08, y + 0.22, 0.45, 0.24, 18, True, accents[i])
        add_text(slide, title, 1.78, y + 0.18, 4.3, 0.24, 17, True, WHITE)
        add_text(slide, desc, 6.0, y + 0.19, 5.8, 0.34, 12, False, TEXT)
        y += 1.26
    footer(slide, 7)

    # 8 отечественное ПО и правовой контур
    slide = prs.slides.add_slide(blank)
    add_bg(slide, None, BG)
    add_text(slide, "Нормативный и технологический контур", 0.68, 0.62, 6.2, 0.34, 28, True, TEXT_DARK)
    add_text(slide, "Для продукта важны не только функции, но и дальнейшая правовая и технологическая упаковка под требования российского рынка.", 0.68, 1.04, 7.8, 0.4, 12.4, False, MUTED_DARK)
    add_rect(slide, 0.76, 1.84, 5.98, 4.92, WHITE, rounded=True, line=LINE)
    add_rect(slide, 6.82, 1.84, 5.74, 4.92, WHITE, rounded=True, line=LINE)
    add_text(slide, "Переход на отечественное ПО", 1.04, 2.18, 3.8, 0.24, 20, True, TEXT_DARK)
    add_bullets(slide, [
        "Запланирован после завершения процесса написания основного кода.",
        "Цель — адаптация продуктового контура под целевой отечественный стек.",
        "Переход должен проводиться без потери логики продукта и UX-подачи."
    ], 1.04, 2.7, 4.95, 1.8, size=14, color=MUTED_DARK)
    add_text(slide, "Правовая регистрация", 7.08, 2.18, 3.2, 0.24, 20, True, TEXT_DARK)
    add_bullets(slide, [
        "Регистрация товарного знака.",
        "Регистрация самого продукта в соответствии с законодательством РФ.",
        "Формирование юридически корректной модели вывода на рынок."
    ], 7.08, 2.7, 4.75, 1.8, size=14, color=MUTED_DARK)
    add_chip(slide, "Roadmap after core-code", 1.04, 5.98, accent=CYAN, width=2.0, light=True)
    add_chip(slide, "Legal readiness", 7.08, 5.98, accent=ORANGE, width=1.5, light=True)
    footer(slide, 8)

    # 9 два продукта
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_split)
    add_rect(slide, 0, 0, SW, SH, NAVY, transparency=0.5)
    add_text(slide, "Два продукта", 0.72, 0.58, 3.2, 0.34, 28, True, WHITE)
    add_text(slide, "Один продукт развивается как корпоративное решение для СТК/ПАО Сбер, второй — как рыночный продукт с доступом к платформе по подписке.", 0.72, 1.0, 8.0, 0.46, 12.5, False, TEXT)
    add_rect(slide, 0.78, 1.82, 5.86, 4.96, NAVY_2, transparency=0.17, rounded=True, line=ORANGE, line_transparency=0.1)
    add_rect(slide, 6.72, 1.82, 5.86, 4.96, NAVY_2, transparency=0.17, rounded=True, line=CYAN, line_transparency=0.1)
    add_chip(slide, "Продукт 1", 1.0, 2.08, accent=ORANGE, width=1.12)
    add_text(slide, "Для ООО «СТК»\n(под Генеральное соглашение\nс ПАО Сбер)", 1.0, 2.45, 4.5, 1.05, 21, True, WHITE)
    add_bullets(slide, [
        "Отдельная разработка под корпоративный контур.",
        "Приведение в соответствие стандартам по кибербезопасности РФ и ПАО Сбер.",
        "Отдельная логика внедрения, эксплуатации и сопровождения."
    ], 1.0, 3.85, 4.75, 2.0, size=14.4)
    add_chip(slide, "Продукт 2", 6.96, 2.08, accent=CYAN, width=1.12)
    add_text(slide, "Для внешнего рынка:\nдоступ к платформе\nкак коммерческому сервису", 6.96, 2.45, 4.45, 1.05, 21, True, WHITE)
    add_bullets(slide, [
        "Продажа доступа к платформе в виде подписки.",
        "Тарифные планы с разным доступом к функционалу.",
        "Абонентская плата как масштабируемая модель монетизации."
    ], 6.96, 3.85, 4.75, 2.0, size=14.4)
    footer(slide, 9)

    # 10 тарифы + финальный титульный-like closure
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_tariffs)
    add_rect(slide, 0, 0, SW, SH, NAVY, transparency=0.52)
    add_text(slide, "Модель тарификации для внешнего рынка", 0.72, 0.58, 6.4, 0.34, 28, True, WHITE)
    add_text(slide, "Предварительная продуктовая логика: разные тарифы дают разный объём функционала и, соответственно, разную стоимость доступа.", 0.72, 1.0, 7.4, 0.42, 12.5, False, TEXT)
    tariffs = [
        ("Start", "от 49 000 ₽/мес", ["базовый пресейл", "стандартные выгрузки", "ограниченный объём доступа"], CYAN),
        ("Pro", "от 119 000 ₽/мес", ["полный расчёт", "AI-модули", "расширенные выходные материалы"], ORANGE),
        ("Enterprise", "индивидуально", ["корпоративная конфигурация", "приоритетная поддержка", "интеграции и SLA"], GREEN),
    ]
    xs = [0.88, 4.48, 8.08]
    for i, (name, price, items, accent) in enumerate(tariffs):
        add_rect(slide, xs[i], 1.86, 3.0, 4.76, WHITE, transparency=0.04, rounded=True, line=accent, line_transparency=0.0)
        add_text(slide, name, xs[i] + 0.2, 2.18, 1.8, 0.24, 20, True, WHITE)
        add_text(slide, price, xs[i] + 0.2, 2.62, 2.1, 0.24, 16, True, accent)
        add_bullets(slide, items, xs[i] + 0.2, 3.28, 2.45, 1.6, size=13.6, color=TEXT)
        add_chip(slide, "Абонентская плата", xs[i] + 0.2, 5.88, accent=accent, width=1.7)
    add_text(slide, "Для версии ООО «СТК» (ПАО Сбер) предполагается отдельная договорная модель поставки вне стандартной тарифной линейки.", 0.9, 6.92, 11.5, 0.22, 10.4, False, MUTED)
    footer(slide, 10)

    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    result = build()
    print(result)
