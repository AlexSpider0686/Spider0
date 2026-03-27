from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


XLSX = "budget_v4.xlsx"
PPT = "k32g_v2_work.pptx"
OUT = "k32g_v2_fixed.pptx"


def fmt_mln(v):
    return f"{v/1_000_000:.1f}".replace(".", ",") + " млн ₽"


def fmt_rub(v):
    return f"{round(v):,}".replace(",", " ") + " ₽"


def set_text(shape, text, size=None, bold=None, color=None, align=None):
    shape.text = text
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        if align is not None:
            p.alignment = align
        for r in p.runs:
            r.font.name = "Arial"
            if size is not None:
                r.font.size = Pt(size)
            if bold is not None:
                r.font.bold = bold
            if color is not None:
                r.font.color.rgb = color


def style_table(shape, rows, font_size=7.6):
    table = shape.table
    for r in range(len(table.rows)):
        left, right = rows[r] if r < len(rows) else ("", "")
        table.cell(r, 0).text = left
        table.cell(r, 1).text = right
        for c in range(2):
            cell = table.cell(r, c)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(font_size)
                    run.font.bold = bool(left) and r == len(rows) - 1
                    run.font.color.rgb = RGBColor(0x22, 0x30, 0x42)


wb = load_workbook(XLSX, data_only=True)
ws = wb[wb.sheetnames[0]]

contract_budget = ws["B9"].value * 1.22
own_budget = ws["B23"].value * 1.22
contract_work = 28_200_000.0
own_work = (ws["B20"].value + ws["B21"].value) * 1.22
area = 60515
contract_m2 = round(contract_budget / area)
own_m2 = round(own_budget / area)

own_rows = [
    ("Материалы", fmt_rub(ws["B16"].value * 1.22)),
    ("Оборудование", fmt_rub(ws["B17"].value * 1.22)),
    ("Инструмент / расходники / СИЗ", fmt_rub(ws["B18"].value * 1.22)),
    ("ФОТ проектирование", fmt_rub(ws["B19"].value * 1.22)),
    ("ФОТ", fmt_rub(ws["B20"].value * 1.22)),
    ("ОПР", fmt_rub(ws["B21"].value * 1.22)),
    ("Работы / ПНР собств. ресурсом", fmt_rub((ws["B20"].value + ws["B21"].value) * 1.22)),
    ("АХР и ОПР", fmt_rub(ws["B22"].value * 1.22)),
    ("Итог бюджета проекта", fmt_rub(ws["B23"].value * 1.22)),
]

contract_rows = [
    ("Материалы", fmt_rub(ws["B3"].value * 1.22)),
    ("Оборудование", fmt_rub(ws["B4"].value * 1.22)),
    ("Базовая стоимость работ", fmt_rub(ws["B5"].value * 1.22)),
    ("Корректировка работ", fmt_rub(ws["C5"].value * 1.22)),
    ("ФОТ проектирование", fmt_rub(ws["B6"].value * 1.22)),
    ("ОПР + АХР и ОПР", fmt_rub((ws["B7"].value + ws["B8"].value) * 1.22)),
    ("Итог бюджета проекта", fmt_rub(ws["B9"].value * 1.22)),
]

prs = Presentation(PPT)

# Common compact budget texts
budget_dual = f"Подрядный: {fmt_mln(contract_budget)}\nСобственный: {fmt_mln(own_budget)}"
works_dual = f"Подрядный: {fmt_mln(contract_work)}\nСобственный: {fmt_mln(own_work)}"

# Slide 1
s = prs.slides[0]
set_text(s.shapes[14], "Бюджет проекта", 8.0, True, RGBColor(0x66, 0x77, 0x89), PP_ALIGN.LEFT)
set_text(s.shapes[15], budget_dual, 7.4, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[17], "Собственный ресурс", 8.0, True, RGBColor(0x66, 0x77, 0x89), PP_ALIGN.LEFT)
set_text(s.shapes[18], fmt_mln(own_budget), 12.2, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[20], "Подрядный ресурс", 8.0, True, RGBColor(0x66, 0x77, 0x89), PP_ALIGN.LEFT)
set_text(s.shapes[21], fmt_mln(contract_budget), 12.2, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)

# Slide 2
s = prs.slides[1]
set_text(s.shapes[8], budget_dual, 9.0, True, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.LEFT)
set_text(s.shapes[17], works_dual, 7.3, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[75], f"Подрядный: {fmt_mln(contract_budget)} / {contract_m2} ₽/м²\nСобственный: {fmt_mln(own_budget)} / {own_m2} ₽/м²", 7.2, False, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.CENTER)
set_text(s.shapes[81], f"Подрядный: {contract_m2} ₽ / м²\nСобственный: {own_m2} ₽ / м²", 8.6, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)

# Slide 4
s = prs.slides[3]
set_text(s.shapes[14], budget_dual, 9.0, True, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.LEFT)
set_text(s.shapes[23], works_dual, 7.2, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(
    s.shapes[28],
    "Подрядный: обор. 46,0% • матер. 14,9% • работы 36,2% • проект 2,8%\n"
    "Собственный: обор. 36,1% • матер. 11,7% • работы 49,9% • проект 2,2%",
    7.1,
    False,
    RGBColor(0x22, 0x30, 0x42),
    PP_ALIGN.LEFT,
)
set_text(s.shapes[50], "486 / 853 ₽ за 1 м кабеля", 8.2, False, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[51], "подрядный / собственный", 7.4, False, RGBColor(0x66, 0x77, 0x89), PP_ALIGN.LEFT)
set_text(s.shapes[53], "12 423 / 21 797 ₽ за 1 устройство", 8.0, False, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[54], "монтаж / адресация", 7.4, False, RGBColor(0x66, 0x77, 0x89), PP_ALIGN.LEFT)

# Slide 6
s = prs.slides[5]
set_text(s.shapes[11], works_dual, 7.2, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[17], budget_dual, 9.0, True, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.LEFT)
set_text(s.shapes[20], fmt_mln(own_budget), 12.5, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[23], fmt_mln(contract_budget), 12.5, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[26], f"{contract_m2} / {own_m2} ₽", 10.4, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[31], f"Удельный бюджет: {own_m2} ₽ / м²", 8.2, True, RGBColor(0x66, 0x77, 0x89), PP_ALIGN.LEFT)
set_text(s.shapes[32], "Работы / ПНР:\nподрядный 28,2 млн ₽\nсобственный 49,5 млн ₽", 6.6, False, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[38], f"Удельный бюджет: {contract_m2} ₽ / м²", 8.2, True, RGBColor(0x66, 0x77, 0x89), PP_ALIGN.LEFT)

style_table(s.shapes[29], own_rows, 7.2)
style_table(s.shapes[36], contract_rows, 7.2)

prs.save(OUT)
print(OUT)
