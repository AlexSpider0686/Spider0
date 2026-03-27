from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt


TARGET = "target_g.pptx"
OUT = "target_g_restyled.pptx"


DARK = RGBColor(0x22, 0x30, 0x42)
SUB = RGBColor(0x66, 0x77, 0x89)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xDC, 0xE7, 0xF3)
GREEN = RGBColor(0x21, 0x98, 0x5D)
NAVY = RGBColor(0x12, 0x34, 0x5B)


def style_text(shape, size=None, bold=None, color=None, align=None, name="Arial", valign=None):
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.word_wrap = True
    if valign is not None:
        tf.vertical_anchor = valign
    for p in tf.paragraphs:
        if align is not None:
            p.alignment = align
        for r in p.runs:
            r.font.name = name
            if size is not None:
                r.font.size = Pt(size)
            if bold is not None:
                r.font.bold = bold
            if color is not None:
                r.font.color.rgb = color


def style_cell(cell, size=8.0, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in tf.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


prs = Presentation(TARGET)

# Apply Arial everywhere first to avoid missing Cyrillic glyphs in edited runs.
for slide in prs.slides:
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            style_text(shape)
        if shape.shape_type == 19:  # table
            for row in shape.table.rows:
                for cell in row.cells:
                    style_cell(cell)

# Slide 1
s = prs.slides[0]
s.shapes[19].text = ""
s.shapes[22].text = ""
style_text(s.shapes[4], 21.5, True, WHITE, PP_ALIGN.LEFT)
style_text(s.shapes[5], 10.0, False, RGBColor(0xD8, 0xE4, 0xF1), PP_ALIGN.LEFT)
for idx in [8, 11, 14, 17, 20]:
    style_text(s.shapes[idx], 8.0, True, SUB, PP_ALIGN.LEFT)
style_text(s.shapes[9], 11.2, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[12], 13.2, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[15], 9.2, True, DARK, PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP)
style_text(s.shapes[18], 13.0, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[21], 13.0, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[23], 8.0, False, SUB, PP_ALIGN.LEFT)
style_text(s.shapes[24], 8.0, False, SUB, PP_ALIGN.CENTER)

# Slide 2
s = prs.slides[1]
style_text(s.shapes[1], 21.5, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[2], 8.8, False, SUB, PP_ALIGN.LEFT)
style_text(s.shapes[7], 12.0, True, PALE, PP_ALIGN.LEFT)
style_text(s.shapes[8], 11.0, True, WHITE, PP_ALIGN.LEFT)
for idx in [10, 13, 16, 19]:
    style_text(s.shapes[idx], 8.0, True, SUB, PP_ALIGN.LEFT)
for idx in [11, 14, 17, 20]:
    style_text(s.shapes[idx], 12.0, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[25], 14.0, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[26], 11.0, False, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[73], 12.0, True, SUB, PP_ALIGN.LEFT)
style_text(s.shapes[75], 10.0, False, DARK, PP_ALIGN.CENTER)
style_text(s.shapes[77], 12.0, True, GREEN, PP_ALIGN.LEFT)
style_text(s.shapes[78], 9.4, False, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[80], 8.0, True, SUB, PP_ALIGN.LEFT)
style_text(s.shapes[81], 10.2, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[82], 8.0, False, SUB, PP_ALIGN.LEFT)
style_text(s.shapes[83], 8.0, False, SUB, PP_ALIGN.CENTER)

# Slide 3
s = prs.slides[2]
style_text(s.shapes[1], 21.5, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[2], 8.8, False, SUB, PP_ALIGN.LEFT)
style_text(s.shapes[4], 11.0, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[5], 10.0, False, DARK, PP_ALIGN.LEFT)

# Slide 4
s = prs.slides[3]
style_text(s.shapes[1], 21.5, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[2], 8.8, False, SUB, PP_ALIGN.LEFT)
style_text(s.shapes[4], 9.9, True, GREEN, PP_ALIGN.LEFT)
style_text(s.shapes[5], 8.7, False, DARK, PP_ALIGN.LEFT)
for idx in [6, 8, 10]:
    style_text(s.shapes[idx], 9.9, True, GREEN, PP_ALIGN.LEFT)
for idx in [7, 9, 11]:
    style_text(s.shapes[idx], 8.7, False, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[13], 12.0, True, PALE, PP_ALIGN.LEFT)
style_text(s.shapes[14], 11.0, True, WHITE, PP_ALIGN.LEFT)
for idx in [16, 19, 22, 25, 28]:
    style_text(s.shapes[idx], 8.0, True, SUB, PP_ALIGN.LEFT)
for idx in [17, 20, 23, 26]:
    style_text(s.shapes[idx], 10.8, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[28], 7.8, False, DARK, PP_ALIGN.LEFT)
for idx in [35, 38, 41, 44, 47, 50, 53]:
    style_text(s.shapes[idx], 9.0, True, SUB, PP_ALIGN.LEFT)
for idx in [36, 39, 42, 45, 48, 51, 54]:
    style_text(s.shapes[idx], 9.2, False, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[58], 8.0, False, SUB, PP_ALIGN.CENTER)
style_text(s.shapes[59], 8.0, False, SUB, PP_ALIGN.CENTER)

# Slide 5
s = prs.slides[4]
style_text(s.shapes[1], 21.5, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[2], 8.8, False, SUB, PP_ALIGN.LEFT)
style_text(s.shapes[4], 11.0, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[7], 11.0, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[27], 8.0, False, SUB, PP_ALIGN.CENTER)

# Slide 6
s = prs.slides[5]
style_text(s.shapes[1], 21.5, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[2], 8.8, False, SUB, PP_ALIGN.LEFT)
for idx in [4, 7, 10, 13, 16, 19, 22, 25]:
    style_text(s.shapes[idx], 8.0, True, SUB if idx not in [16] else PALE, PP_ALIGN.LEFT)
for idx in [5, 8, 11, 14]:
    style_text(s.shapes[idx], 12.0, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[17], 11.0, True, WHITE, PP_ALIGN.LEFT)
style_text(s.shapes[20], 12.5, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[23], 12.5, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[26], 12.0, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[28], 11.5, True, GREEN, PP_ALIGN.LEFT)
style_text(s.shapes[31], 10.0, True, SUB, PP_ALIGN.LEFT)
style_text(s.shapes[32], 10.0, False, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[35], 11.5, True, GREEN, PP_ALIGN.LEFT)
style_text(s.shapes[38], 10.0, True, SUB, PP_ALIGN.LEFT)
style_text(s.shapes[39], 8.6, False, SUB, PP_ALIGN.LEFT)
style_text(s.shapes[41], 8.0, False, SUB, PP_ALIGN.CENTER)
style_text(s.shapes[42], 8.0, False, SUB, PP_ALIGN.CENTER)

for table_idx in [29, 36]:
    table = s.shapes[table_idx].table
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            style_cell(table.cell(r, c), 8.1, bold=(r == 0), align=PP_ALIGN.LEFT)

# Slide 7
s = prs.slides[6]
style_text(s.shapes[1], 21.5, True, NAVY, PP_ALIGN.LEFT)
style_text(s.shapes[2], 14.0, False, SUB, PP_ALIGN.LEFT)
for idx in [6, 8, 10, 12, 14, 16, 18]:
    style_text(s.shapes[idx], 12.5, True, NAVY, PP_ALIGN.CENTER)
for idx in [19, 21, 23, 25, 27, 29]:
    style_text(s.shapes[idx], 14.0, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[32], 11.5, True, WHITE, PP_ALIGN.CENTER)
style_text(s.shapes[34], 14.0, True, NAVY, PP_ALIGN.LEFT)
style_text(s.shapes[35], 10.0, False, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[38], 14.0, True, NAVY, PP_ALIGN.LEFT)
style_text(s.shapes[39], 10.0, False, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[40], 8.0, False, SUB, PP_ALIGN.CENTER)
style_text(s.shapes[41], 8.0, False, SUB, PP_ALIGN.CENTER)

# Slide 8
s = prs.slides[7]
style_text(s.shapes[1], 21.5, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[2], 8.8, False, SUB, PP_ALIGN.LEFT)
style_text(s.shapes[6], 17.0, True, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[7], 11.2, False, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[8], 17.0, False, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[9], 12.0, False, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[10], 17.0, False, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[11], 12.0, False, DARK, PP_ALIGN.LEFT)
style_text(s.shapes[15], 8.0, False, SUB, PP_ALIGN.CENTER)
style_text(s.shapes[16], 8.0, False, SUB, PP_ALIGN.CENTER)

prs.save(OUT)
print(OUT)
