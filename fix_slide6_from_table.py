from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


XLSX = "budget_v4.xlsx"
PPT = "k32g_v2_work.pptx"
OUT = "k32g_v2_slide6_fixed.pptx"


def rub(v):
    return f"{round(v):,}".replace(",", " ") + " ₽"


def mln(v):
    return f"{v/1_000_000:.1f}".replace(".", ",") + " млн ₽"


def set_text(shape, text, size=None, bold=None, color=None, align=None):
    shape.text = text
    if not getattr(shape, "has_text_frame", False):
        return
    for p in shape.text_frame.paragraphs:
        if align is not None:
            p.alignment = align
        for r in p.runs:
            r.font.name = "Arial"
            if size is not None:
                r.font.size = Pt(size)
            if bold is not None:
                r.font.bold = bold
            if color is not None:
                r.font.color.rgb = color


def fill_table(shape, rows, font_size=7.0):
    table = shape.table
    for r in range(len(table.rows)):
        left, right = rows[r] if r < len(rows) else ("", "")
        table.cell(r, 0).text = left
        table.cell(r, 1).text = right
        for c in range(2):
            cell = table.cell(r, c)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(font_size)
                    run.font.bold = left == "Итог бюджета проекта"
                    run.font.color.rgb = RGBColor(0x22, 0x30, 0x42)


wb = load_workbook(XLSX, data_only=True)
ws = wb[wb.sheetnames[0]]

# Final totals remain unchanged per user.
own_total = ws["F9"].value
contract_total = ws["F23"].value

# Values exactly from table logic requested by user.
mat_rev = ws["F3"].value
eq_rev = ws["F4"].value

contract_opr_cost = ws["B7"].value
contract_ahr_cost = ws["B8"].value

own_instr_cost = ws["B18"].value
own_opr_cost = ws["B21"].value
own_ahr_cost = ws["B22"].value

prs = Presentation(PPT)

# Update upper captions on slides 1/2/4/6 so they match requested totals.
budget_dual = f"Собственный: {mln(own_total)}\nПодрядный: {mln(contract_total)}"

for slide_no, own_idx, contract_idx, budget_idx in [
    (1, 19, 22, 16),
    (6, 21, 24, 18),
]:
    slide = prs.slides[slide_no - 1]
    set_text(slide.shapes[budget_idx - 1], budget_dual, 8.8 if slide_no == 6 else 7.3, True, RGBColor(0x22, 0x30, 0x42) if slide_no == 1 else RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.LEFT)
    set_text(slide.shapes[own_idx - 1], mln(own_total), 12.2, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
    set_text(slide.shapes[contract_idx - 1], mln(contract_total), 12.2, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)

for slide_no, budget_idx in [(2, 9), (4, 15)]:
    slide = prs.slides[slide_no - 1]
    set_text(slide.shapes[budget_idx - 1], budget_dual, 8.8, True, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.LEFT)

# Slide 6 detailed budgets
s = prs.slides[5]
set_text(s.shapes[5], mln(eq_rev), 12.0, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[8], mln(mat_rev), 12.0, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[17], budget_dual, 8.8, True, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.LEFT)
set_text(s.shapes[20], mln(own_total), 12.2, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[23], mln(contract_total), 12.2, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)

own_rows = [
    ("Материалы (выручка)", rub(mat_rev)),
    ("Оборудование (выручка)", rub(eq_rev)),
    ("Инстр. / расходники / СИЗ", rub(own_instr_cost)),
    ("ОПР (себестоимость)", rub(own_opr_cost)),
    ("АХР и ОПР (себестоимость)", rub(own_ahr_cost)),
    ("Работы / ПНР", "28 200 000 ₽"),
    ("Итог бюджета проекта", rub(own_total)),
]
contract_rows = [
    ("Материалы (выручка)", rub(mat_rev)),
    ("Оборудование (выручка)", rub(eq_rev)),
    ("ОПР (себестоимость)", rub(contract_opr_cost)),
    ("АХР и ОПР (себестоимость)", rub(contract_ahr_cost)),
    ("Инструменты / расходники", "—"),
    ("Работы / ПНР", "49 479 558 ₽"),
    ("Итог бюджета проекта", rub(contract_total)),
]
fill_table(s.shapes[29], own_rows, 6.8)
fill_table(s.shapes[36], contract_rows, 6.8)

prs.save(OUT)
print(OUT)
