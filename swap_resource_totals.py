from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


PPT = "k32g_v2_work.pptx"
OUT = "k32g_v2_swapped.pptx"

OWN_BUDGET = 106_526_532.42
CONTRACT_BUDGET = 143_108_099.38
OWN_WORK = 28_200_000.00
CONTRACT_WORK = 49_479_557.83
OWN_M2 = 1760
CONTRACT_M2 = 2365


def set_text(shape, text, size=None, bold=None, color=None, align=None):
    shape.text = text
    if not getattr(shape, "has_text_frame", False):
        return
    for p in shape.text_frame.paragraphs:
        if align is not None:
            p.alignment = align
        for r in p.runs:
            r.font.name = "Arial"
            if size is not None:
                r.font.size = Pt(size)
            if bold is not None:
                r.font.bold = bold
            if color is not None:
                r.font.color.rgb = color


def fill_table(shape, rows, size=7.2):
    table = shape.table
    for r in range(len(table.rows)):
        left, right = rows[r] if r < len(rows) else ("", "")
        table.cell(r, 0).text = left
        table.cell(r, 1).text = right
        for c in range(2):
            cell = table.cell(r, c)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(size)
                    run.font.bold = left == "Итог бюджета проекта"
                    run.font.color.rgb = RGBColor(0x22, 0x30, 0x42)


prs = Presentation(PPT)

budget_dual = "Собственный: 106,5 млн ₽\nПодрядный: 143,1 млн ₽"
works_dual = "Собственный: 28,2 млн ₽\nПодрядный: 49,5 млн ₽"

# Slide 1
s = prs.slides[0]
set_text(s.shapes[15], budget_dual, 7.3, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[18], "106,5 млн ₽", 12.2, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[21], "143,1 млн ₽", 12.2, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)

# Slide 2
s = prs.slides[1]
set_text(s.shapes[8], budget_dual, 8.8, True, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.LEFT)
set_text(s.shapes[17], works_dual, 7.2, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(
    s.shapes[75],
    "Собственный: 106,5 млн ₽ / 1 760 ₽/м²\nПодрядный: 143,1 млн ₽ / 2 365 ₽/м²",
    7.0,
    False,
    RGBColor(0x22, 0x30, 0x42),
    PP_ALIGN.CENTER,
)
set_text(
    s.shapes[81],
    "Собственный: 1 760 ₽ / м²\nПодрядный: 2 365 ₽ / м²",
    8.4,
    True,
    RGBColor(0x22, 0x30, 0x42),
    PP_ALIGN.LEFT,
)

# Slide 4
s = prs.slides[3]
set_text(s.shapes[14], budget_dual, 8.8, True, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.LEFT)
set_text(s.shapes[23], works_dual, 7.2, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(
    s.shapes[28],
    "Собственный: обор. 46,0% • матер. 14,9% • работы 36,2% • проект 2,8%\n"
    "Подрядный: обор. 36,1% • матер. 11,7% • работы 49,9% • проект 2,2%",
    7.0,
    False,
    RGBColor(0x22, 0x30, 0x42),
    PP_ALIGN.LEFT,
)
set_text(s.shapes[50], "486 / 853 ₽ за 1 м кабеля", 8.0, False, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[51], "собственный / подрядный", 7.2, False, RGBColor(0x66, 0x77, 0x89), PP_ALIGN.LEFT)
set_text(s.shapes[53], "12 423 / 21 797 ₽ за 1 устройство", 7.8, False, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)

# Slide 6
s = prs.slides[5]
set_text(s.shapes[11], works_dual, 7.1, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[17], budget_dual, 8.8, True, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.LEFT)
set_text(s.shapes[20], "106,5 млн ₽", 12.2, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[23], "143,1 млн ₽", 12.2, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[26], "1 760 / 2 365 ₽", 9.8, True, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[31], "Удельный бюджет: 1 760 ₽ / м²", 8.0, True, RGBColor(0x66, 0x77, 0x89), PP_ALIGN.LEFT)
set_text(s.shapes[32], "Работы / ПНР:\nсобственный 28,2 млн ₽\nподрядный 49,5 млн ₽", 6.5, False, RGBColor(0x22, 0x30, 0x42), PP_ALIGN.LEFT)
set_text(s.shapes[38], "Удельный бюджет: 2 365 ₽ / м²", 8.0, True, RGBColor(0x66, 0x77, 0x89), PP_ALIGN.LEFT)

own_table_rows = [
    ("Материалы", "11 600 000 ₽"),
    ("Оборудование", "35 800 000 ₽"),
    ("Базовая стоимость работ", "28 200 000 ₽"),
    ("Корректировка работ", "13 553 625 ₽"),
    ("ФОТ проектирование", "2 209 864 ₽"),
    ("ОПР + АХР и ОПР", "2 677 592 ₽"),
    ("Итог бюджета проекта", "106 526 532 ₽"),
]
contract_table_rows = [
    ("Материалы", "11 600 000 ₽"),
    ("Оборудование", "35 800 000 ₽"),
    ("ФОТ проектирование", "2 209 864 ₽"),
    ("ФОТ + ОПР", "85 735 945 ₽"),
    ("Инструмент / расходники / СИЗ", "4 834 185 ₽"),
    ("АХР и ОПР", "23 777 134 ₽"),
    ("Итог бюджета проекта", "143 108 099 ₽"),
]
fill_table(s.shapes[29], own_table_rows, 7.0)
fill_table(s.shapes[36], contract_table_rows, 7.0)

prs.save(OUT)
print(OUT)
